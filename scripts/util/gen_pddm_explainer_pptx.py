#!/usr/bin/env python3
"""
Generate a simple PPTX explaining the Practice Depth Determination Model (PDDM)
— how teacher orientation survey responses map to FP + Depth, and how
dominant FP and average depth are derived per teacher.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

OUT = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(
    os.path.abspath(__file__)))), "output", "PDDM_Explained_Simple.pptx")

# Colors
NAVY = RGBColor(0x1E, 0x3A, 0x5F)
BLUE = RGBColor(0x25, 0x63, 0xEB)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
AMBER = RGBColor(0xD9, 0x77, 0x06)
RED = RGBColor(0xDC, 0x26, 0x26)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
TEAL = RGBColor(0x0D, 0x94, 0x88)
GRAY = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x2E)
LIGHT_BG = RGBColor(0xF5, 0xF6, 0xFA)
LIGHT_BLUE = RGBColor(0xDB, 0xEA, 0xFE)
LIGHT_GREEN = RGBColor(0xDC, 0xFC, 0xE7)
LIGHT_AMBER = RGBColor(0xFE, 0xF3, 0xC7)
LIGHT_RED = RGBColor(0xFE, 0xE2, 0xE2)
LIGHT_PURPLE = RGBColor(0xED, 0xE9, 0xFE)
LIGHT_TEAL = RGBColor(0xCC, 0xFB, 0xF1)

D1_COLOR = LIGHT_RED
D2_COLOR = LIGHT_AMBER
D3_COLOR = LIGHT_BLUE
D4_COLOR = LIGHT_GREEN


def shape(slide, left, top, w, h, fill, text="", sz=12, fc=BLACK, bold=False, align=PP_ALIGN.LEFT):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill; s.line.fill.background()
    if text:
        tf = s.text_frame; tf.word_wrap = True
        tf.margin_left = Pt(8); tf.margin_right = Pt(8); tf.margin_top = Pt(4)
        p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(sz)
        p.font.color.rgb = fc; p.font.bold = bold; p.alignment = align
    return s


def txt(slide, left, top, w, h, text, sz=14, fc=BLACK, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(sz)
    p.font.color.rgb = fc; p.font.bold = bold; p.alignment = align
    return tb


def multi_txt(slide, left, top, w, h, lines, sz=12, fc=BLACK):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(sz); p.font.color.rgb = fc; p.space_after = Pt(3)
    return tb


def slide_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    shape(s, Inches(0), Inches(0), Inches(10), Inches(7.5), NAVY)
    txt(s, Inches(0.8), Inches(1.2), Inches(8.4), Inches(1),
        "PDDM", 44, WHITE, True, PP_ALIGN.CENTER)
    txt(s, Inches(0.8), Inches(2.2), Inches(8.4), Inches(0.5),
        "Practice Depth Determination Model", 24, RGBColor(0xBF, 0xDB, 0xFE), False, PP_ALIGN.CENTER)
    txt(s, Inches(0.8), Inches(3.5), Inches(8.4), Inches(1),
        "How we determine WHERE a teacher's attention goes (FP)\nand HOW DEEP their practice is (D1-D4)\nfrom their survey responses",
        16, RGBColor(0x93, 0xC5, 0xFD), False, PP_ALIGN.CENTER)
    txt(s, Inches(0.8), Inches(5.5), Inches(8.4), Inches(0.5),
        "Instruments: 1234 (Teacher EN) | 103 (Teacher MR) | 7890 (Leader EN) | 104 (Leader MR)",
        13, RGBColor(0x93, 0xC5, 0xFD), False, PP_ALIGN.CENTER)
    txt(s, Inches(0.8), Inches(6.5), Inches(8.4), Inches(0.5),
        "Project Kshitij | Myelin | Confidential", 11, RGBColor(0x60, 0x7D, 0x8B), False, PP_ALIGN.CENTER)


def slide_big_picture(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    shape(s, Inches(0), Inches(0), Inches(10), Inches(0.9), NAVY)
    txt(s, Inches(0.5), Inches(0.15), Inches(9), Inches(0.6), "The Big Picture", 28, WHITE, True)

    txt(s, Inches(0.5), Inches(1.2), Inches(9), Inches(0.5),
        "Every question in the orientation survey has TWO hidden tags on each answer option:", 14, NAVY, True)

    # Two tags visual
    shape(s, Inches(1), Inches(2.0), Inches(3.5), Inches(1.8), LIGHT_BLUE)
    txt(s, Inches(1.2), Inches(2.1), Inches(3.2), Inches(0.4), "Tag 1: FP (Focus Point)", 16, BLUE, True)
    multi_txt(s, Inches(1.2), Inches(2.6), Inches(3.2), Inches(1.1), [
        "Which NEP philosophy does",
        "this answer align with?",
        "",
        "FP1, FP2, FP3, FP4, or FP5"
    ], 12, GRAY)

    shape(s, Inches(5.5), Inches(2.0), Inches(3.5), Inches(1.8), LIGHT_GREEN)
    txt(s, Inches(5.7), Inches(2.1), Inches(3.2), Inches(0.4), "Tag 2: DEPTH", 16, GREEN, True)
    multi_txt(s, Inches(5.7), Inches(2.6), Inches(3.2), Inches(1.1), [
        "How deep is this practice?",
        "",
        "D1 (surface) → D2 (responsive)",
        "→ D3 (diagnostic) → D4 (adaptive)"
    ], 12, GRAY)

    # Arrow
    txt(s, Inches(0.5), Inches(4.2), Inches(9), Inches(0.3), "▼", 24, GRAY, True, PP_ALIGN.CENTER)

    shape(s, Inches(1), Inches(4.7), Inches(8), Inches(1.6), LIGHT_PURPLE)
    txt(s, Inches(1.2), Inches(4.8), Inches(7.6), Inches(0.4), "From responses, we derive per teacher:", 14, PURPLE, True)
    multi_txt(s, Inches(1.2), Inches(5.3), Inches(7.6), Inches(0.9), [
        "1. Dominant FP — which philosophy do they naturally lean towards?",
        "2. Average Depth — how deep is their practice across all FPs?",
        "3. Per-FP Depth — how deep is practice specifically for FP1, FP2, etc.?",
    ], 13, BLACK)

    shape(s, Inches(0.5), Inches(6.6), Inches(9), Inches(0.6), LIGHT_AMBER)
    txt(s, Inches(0.7), Inches(6.65), Inches(8.6), Inches(0.5),
        "The teacher doesn't know about FP or Depth tags. They just pick the answer that feels right.\nThe tags are invisible — built into the instrument design.",
        10, AMBER)


def slide_fp_explained(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    shape(s, Inches(0), Inches(0), Inches(10), Inches(0.9), BLUE)
    txt(s, Inches(0.5), Inches(0.15), Inches(9), Inches(0.6),
        "The 5 Foundational Philosophies (FPs)", 28, WHITE, True)

    txt(s, Inches(0.5), Inches(1.1), Inches(9), Inches(0.3),
        "FP = where a teacher's attention naturally goes first. It is INTENT, not skill.", 13, GRAY)

    fps = [
        ("FP1", "Each Child is Unique", "Notices learner differences, adapts to individuals",
         "Q1-Q3 in orientation survey", LIGHT_BLUE, BLUE),
        ("FP2", "Holistic & Experiential Learning", "Connects learning to real life, goes beyond textbook",
         "Q7-Q9 in orientation survey", LIGHT_GREEN, GREEN),
        ("FP3", "Reflective Practitioner", "Reflects on practice, seeks feedback, experiments",
         "Q10-Q12 in orientation survey", LIGHT_PURPLE, PURPLE),
        ("FP4", "Assessment for Learning", "Uses assessment to diagnose, not just rank",
         "Q4-Q6 in orientation survey", LIGHT_AMBER, AMBER),
        ("FP5", "Collaboration & Community", "Works with parents, peers, community as partners",
         "Q13-Q15 in orientation survey", LIGHT_TEAL, TEAL),
    ]

    y = 1.6
    for code, name, meaning, source, bg, fg in fps:
        shape(s, Inches(0.5), Inches(y), Inches(9), Inches(0.9), bg)
        shape(s, Inches(0.5), Inches(y), Inches(0.8), Inches(0.9), fg)
        txt(s, Inches(0.55), Inches(y + 0.2), Inches(0.7), Inches(0.4), code, 16, WHITE, True, PP_ALIGN.CENTER)
        txt(s, Inches(1.5), Inches(y + 0.02), Inches(4), Inches(0.35), name, 14, BLACK, True)
        txt(s, Inches(1.5), Inches(y + 0.38), Inches(4.5), Inches(0.3), meaning, 11, GRAY)
        txt(s, Inches(6.2), Inches(y + 0.25), Inches(3.2), Inches(0.3), source, 10, fg, False, PP_ALIGN.RIGHT)
        y += 1.0


def slide_depth_explained(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    shape(s, Inches(0), Inches(0), Inches(10), Inches(0.9), GREEN)
    txt(s, Inches(0.5), Inches(0.15), Inches(9), Inches(0.6),
        "The 4 Practice Depth Levels (D1-D4)", 28, WHITE, True)

    txt(s, Inches(0.5), Inches(1.1), Inches(9), Inches(0.3),
        "Depth = how far intent translates into classroom action. Not good/bad — it's a progression.", 13, GRAY)

    depths = [
        ("D1", "Procedural", "Follows standard procedure", "\"I teach the chapter as planned\"",
         "Rule-based, default responses. Does what's expected.", D1_COLOR, RED),
        ("D2", "Responsive", "Adjusts based on the situation", "\"I noticed students struggling, so I slowed down\"",
         "Reacts to what happens, but doesn't investigate why.", D2_COLOR, AMBER),
        ("D3", "Diagnostic", "Uses evidence to understand WHY", "\"I gave a quick check to find the exact misconception\"",
         "Investigates root causes. Uses data to decide next steps.", D3_COLOR, BLUE),
        ("D4", "Adaptive", "Intentionally redesigns practice", "\"I created a new activity based on what diagnostics showed\"",
         "Creates new approaches. Practice evolves based on evidence.", D4_COLOR, GREEN),
    ]

    y = 1.7
    for code, label, short, quote, full, bg, fg in depths:
        shape(s, Inches(0.5), Inches(y), Inches(9), Inches(1.15), bg)
        shape(s, Inches(0.5), Inches(y), Inches(0.8), Inches(1.15), fg)
        txt(s, Inches(0.55), Inches(y + 0.3), Inches(0.7), Inches(0.4), code, 20, WHITE, True, PP_ALIGN.CENTER)
        txt(s, Inches(1.5), Inches(y + 0.02), Inches(3), Inches(0.35), label, 16, BLACK, True)
        txt(s, Inches(4.5), Inches(y + 0.02), Inches(5), Inches(0.35), short, 12, fg, True)
        txt(s, Inches(1.5), Inches(y + 0.4), Inches(7.8), Inches(0.3), quote, 11, PURPLE, False)
        txt(s, Inches(1.5), Inches(y + 0.72), Inches(7.8), Inches(0.35), full, 10, GRAY)
        y += 1.3

    shape(s, Inches(0.5), Inches(6.9), Inches(9), Inches(0.4), LIGHT_AMBER)
    txt(s, Inches(0.7), Inches(6.93), Inches(8.6), Inches(0.35),
        "D1 → D2 → D3 → D4 is a progression, not a judgement. Most teachers are at D1-D2. D3+ is the goal.", 10, AMBER)


def slide_how_tagging_works(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    shape(s, Inches(0), Inches(0), Inches(10), Inches(0.9), NAVY)
    txt(s, Inches(0.5), Inches(0.15), Inches(9), Inches(0.6),
        "How It Works: One Question, Four Tags", 26, WHITE, True)

    txt(s, Inches(0.5), Inches(1.1), Inches(9), Inches(0.4),
        "Example: Q1. \"When two students learn at very different speeds, what do you usually do?\"", 13, NAVY, True)

    options = [
        ("A", "Continue teaching at the same pace for everyone", "FP2", "D1", "Procedural — ignores difference", D1_COLOR),
        ("B", "Give extra homework to the slower learner", "FP1", "D1", "Notices difference but surface response", D1_COLOR),
        ("C", "Provide varied support so both can progress", "FP1", "D3", "Diagnostic — intentionally adapts", D3_COLOR),
        ("D", "Ask the faster learner to help the slower one", "FP5", "D2", "Responsive — uses peer support", D2_COLOR),
    ]

    # Header
    y = 1.8
    for x, label, w in [(0.5, "Option", 0.7), (1.3, "Answer Text", 4.2), (5.6, "FP Tag", 0.9), (6.6, "Depth", 0.8), (7.5, "What It Reveals", 2.2)]:
        shape(s, Inches(x), Inches(y), Inches(w), Inches(0.35), NAVY)
        txt(s, Inches(x + 0.05), Inches(y), Inches(w), Inches(0.3), label, 10, WHITE, True)
    y += 0.4

    for opt, text, fp, depth, reading, bg in options:
        shape(s, Inches(0.5), Inches(y), Inches(0.7), Inches(0.55), bg)
        txt(s, Inches(0.55), Inches(y + 0.1), Inches(0.6), Inches(0.3), opt, 14, BLACK, True, PP_ALIGN.CENTER)
        shape(s, Inches(1.3), Inches(y), Inches(4.2), Inches(0.55), LIGHT_BG)
        txt(s, Inches(1.35), Inches(y + 0.05), Inches(4.1), Inches(0.45), text, 10, BLACK)
        shape(s, Inches(5.6), Inches(y), Inches(0.9), Inches(0.55), LIGHT_BLUE)
        txt(s, Inches(5.6), Inches(y + 0.1), Inches(0.9), Inches(0.3), fp, 12, BLUE, True, PP_ALIGN.CENTER)
        depth_bg = {D1_COLOR: D1_COLOR, D2_COLOR: D2_COLOR, D3_COLOR: D3_COLOR}
        shape(s, Inches(6.6), Inches(y), Inches(0.8), Inches(0.55), bg)
        txt(s, Inches(6.6), Inches(y + 0.1), Inches(0.8), Inches(0.3), depth, 12, BLACK, True, PP_ALIGN.CENTER)
        shape(s, Inches(7.5), Inches(y), Inches(2.2), Inches(0.55), LIGHT_BG)
        txt(s, Inches(7.55), Inches(y + 0.05), Inches(2.1), Inches(0.45), reading, 9, GRAY)
        y += 0.6

    # Explanation
    txt(s, Inches(0.5), Inches(4.5), Inches(9), Inches(0.3),
        "If this teacher picks option C:", 14, NAVY, True)
    shape(s, Inches(0.5), Inches(4.9), Inches(9), Inches(0.8), LIGHT_GREEN)
    multi_txt(s, Inches(0.7), Inches(4.95), Inches(8.6), Inches(0.7), [
        "  FP1 gets +1 vote (this teacher attends to learner uniqueness)",
        "  D3 gets +1 depth point (this teacher uses diagnostic-level practice)",
        "  This single answer adds evidence to both the FP and depth calculations."
    ], 12, BLACK)

    # Rubric levels
    txt(s, Inches(0.5), Inches(5.9), Inches(9), Inches(0.3),
        "3-Level Rubric (invisible to teacher):", 13, NAVY, True)
    for x, level, label, desc, bg, fg in [
        (0.5, "Level 1", "NEP-Aligned", "D3 responses — diagnostic, intentional", LIGHT_GREEN, GREEN),
        (3.5, "Level 2", "Partially Aligned", "D2 responses — responsive, situational", LIGHT_AMBER, AMBER),
        (6.5, "Level 3", "Traditional", "D1 responses — procedural, compliance", LIGHT_RED, RED),
    ]:
        shape(s, Inches(x), Inches(6.3), Inches(2.8), Inches(0.9), bg)
        txt(s, Inches(x + 0.1), Inches(6.33), Inches(2.6), Inches(0.3), level + ": " + label, 11, fg, True, PP_ALIGN.CENTER)
        txt(s, Inches(x + 0.1), Inches(6.63), Inches(2.6), Inches(0.3), desc, 9, GRAY, False, PP_ALIGN.CENTER)


def slide_computation(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    shape(s, Inches(0), Inches(0), Inches(10), Inches(0.9), PURPLE)
    txt(s, Inches(0.5), Inches(0.15), Inches(9), Inches(0.6),
        "How Dominant FP & Depth Are Computed", 26, WHITE, True)

    txt(s, Inches(0.5), Inches(1.1), Inches(9), Inches(0.3),
        "A teacher answers ~20 questions. Each answer has an FP tag and a Depth tag.", 13, GRAY)

    # Step 1
    shape(s, Inches(0.3), Inches(1.6), Inches(4.5), Inches(2.8), LIGHT_BLUE)
    txt(s, Inches(0.5), Inches(1.65), Inches(4.2), Inches(0.35), "Step 1: Count FP Votes", 15, BLUE, True)
    multi_txt(s, Inches(0.5), Inches(2.1), Inches(4.2), Inches(2.0), [
        "Example teacher's 20 answers:",
        "  FP1 appeared 6 times",
        "  FP2 appeared 4 times",
        "  FP3 appeared 3 times",
        "  FP4 appeared 5 times",
        "  FP5 appeared 2 times",
        "",
        "Dominant FP = FP1 (highest count)",
        "This teacher's attention goes to",
        "\"Each Child is Unique\" first."
    ], 11, BLACK)

    # Step 2
    shape(s, Inches(5.2), Inches(1.6), Inches(4.5), Inches(2.8), LIGHT_GREEN)
    txt(s, Inches(5.4), Inches(1.65), Inches(4.2), Inches(0.35), "Step 2: Average Depth", 15, GREEN, True)
    multi_txt(s, Inches(5.4), Inches(2.1), Inches(4.2), Inches(2.0), [
        "Same teacher's depth tags:",
        "  D1 appeared 8 times (score: 1 each)",
        "  D2 appeared 7 times (score: 2 each)",
        "  D3 appeared 5 times (score: 3 each)",
        "  D4 appeared 0 times (score: 4 each)",
        "",
        "Average = (8×1 + 7×2 + 5×3) / 20",
        "        = (8 + 14 + 15) / 20 = 1.85",
        "",
        "Depth Band = D1-D2 (Low: < 2.2)"
    ], 11, BLACK)

    # Step 3
    shape(s, Inches(0.3), Inches(4.7), Inches(9.4), Inches(1.3), LIGHT_PURPLE)
    txt(s, Inches(0.5), Inches(4.75), Inches(9), Inches(0.35), "Step 3: Per-FP Depth (for SET analysis)", 15, PURPLE, True)
    multi_txt(s, Inches(0.5), Inches(5.2), Inches(9), Inches(0.7), [
        "For each FP, what's the average depth of responses tagged with that FP?",
        "  FP1 depth = avg of depths when teacher chose FP1 answers → e.g., 2.3 (just above D2)",
        "  FP4 depth = avg of depths when teacher chose FP4 answers → e.g., 1.4 (solidly D1)",
        "This tells SET: \"This teacher INTENDS assessment (FP4) but only at procedural level (D1).\""
    ], 11, BLACK)

    # Summary box
    shape(s, Inches(0.3), Inches(6.2), Inches(9.4), Inches(1.0), NAVY)
    txt(s, Inches(0.5), Inches(6.25), Inches(9), Inches(0.3), "Per Teacher Output:", 14, WHITE, True)
    multi_txt(s, Inches(0.5), Inches(6.55), Inches(9), Inches(0.6), [
        "Dominant FP: FP1  |  Average Depth: 1.85 (D1-D2)  |  FP1 Depth: 2.3  |  FP4 Depth: 1.4",
        "→ This feeds into SET1 (orientation × depth alignment) at school/cluster level"
    ], 11, RGBColor(0xBF, 0xDB, 0xFE))


def slide_full_example(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    shape(s, Inches(0), Inches(0), Inches(10), Inches(0.9), NAVY)
    txt(s, Inches(0.5), Inches(0.15), Inches(9), Inches(0.6),
        "Full Example: From Survey to Insight", 26, WHITE, True)

    # Teacher answers
    txt(s, Inches(0.3), Inches(1.1), Inches(9.4), Inches(0.3),
        "Ms. Sandhya answers 15 orientation questions (setCode: 1234):", 13, NAVY, True)

    shape(s, Inches(0.3), Inches(1.5), Inches(9.4), Inches(1.8), LIGHT_BG)
    qs = [
        ("Q1", "Varied support for both", "FP1", "D3"),
        ("Q2", "Observing classroom behaviour", "FP1", "D3"),
        ("Q3", "Low-pressure activities", "FP1", "D3"),
        ("Q4", "Quick informal check", "FP4", "D2"),
        ("Q5", "Diagnostic tasks for gaps", "FP4", "D3"),
        ("Q6", "Weekly informal checks", "FP4", "D2"),
        ("Q7", "Building competencies", "FP2", "D3"),
        ("Q8", "Real-life application tasks", "FP2", "D3"),
    ]
    qs2 = [
        ("Q9", "Real-life problem solving", "FP2", "D3"),
        ("Q10", "Last few months", "FP3", "D2"),
        ("Q11", "Helps me improve", "FP3", "D2"),
        ("Q12", "Proactively learn", "FP3", "D3"),
        ("Q13", "Monthly updates", "FP5", "D2"),
        ("Q14", "Strengths & areas", "FP1", "D2"),
        ("Q15", "Share activities", "FP5", "D2"),
    ]

    x, y = 0.4, 1.55
    for q, answer, fp, d in qs:
        txt(s, Inches(x), Inches(y), Inches(1.1), Inches(0.22), f"{q}: {answer}", 7, BLACK)
        bg = {D1_COLOR: D1_COLOR, "D1": D1_COLOR, "D2": D2_COLOR, "D3": D3_COLOR, "D4": D4_COLOR}[d]
        shape(s, Inches(x + 3.5), Inches(y), Inches(0.5), Inches(0.2), LIGHT_BLUE)
        txt(s, Inches(x + 3.5), Inches(y - 0.02), Inches(0.5), Inches(0.2), fp, 7, BLUE, True, PP_ALIGN.CENTER)
        shape(s, Inches(x + 4.05), Inches(y), Inches(0.4), Inches(0.2), bg)
        txt(s, Inches(x + 4.05), Inches(y - 0.02), Inches(0.4), Inches(0.2), d, 7, BLACK, True, PP_ALIGN.CENTER)
        y += 0.22

    x, y = 5.1, 1.55
    for q, answer, fp, d in qs2:
        txt(s, Inches(x), Inches(y), Inches(1.1), Inches(0.22), f"{q}: {answer}", 7, BLACK)
        bg = {"D1": D1_COLOR, "D2": D2_COLOR, "D3": D3_COLOR, "D4": D4_COLOR}[d]
        shape(s, Inches(x + 3.5), Inches(y), Inches(0.5), Inches(0.2), LIGHT_BLUE)
        txt(s, Inches(x + 3.5), Inches(y - 0.02), Inches(0.5), Inches(0.2), fp, 7, BLUE, True, PP_ALIGN.CENTER)
        shape(s, Inches(x + 4.05), Inches(y), Inches(0.4), Inches(0.2), bg)
        txt(s, Inches(x + 4.05), Inches(y - 0.02), Inches(0.4), Inches(0.2), d, 7, BLACK, True, PP_ALIGN.CENTER)
        y += 0.22

    # Computation
    shape(s, Inches(0.3), Inches(3.5), Inches(4.5), Inches(1.5), LIGHT_BLUE)
    txt(s, Inches(0.5), Inches(3.55), Inches(4.2), Inches(0.3), "FP Vote Count:", 13, BLUE, True)
    multi_txt(s, Inches(0.5), Inches(3.9), Inches(4.2), Inches(1.0), [
        "FP1: 4 votes  |  FP2: 3 votes  |  FP3: 3 votes",
        "FP4: 3 votes  |  FP5: 2 votes",
        "",
        "Dominant FP = FP1 (Each Child is Unique)"
    ], 11, BLACK)

    shape(s, Inches(5.2), Inches(3.5), Inches(4.5), Inches(1.5), LIGHT_GREEN)
    txt(s, Inches(5.4), Inches(3.55), Inches(4.2), Inches(0.3), "Depth Calculation:", 13, GREEN, True)
    multi_txt(s, Inches(5.4), Inches(3.9), Inches(4.2), Inches(1.0), [
        "D2: 6 responses  |  D3: 9 responses",
        "Avg = (6×2 + 9×3) / 15 = 39/15 = 2.6",
        "",
        "Depth Band = MID (2.2 – 3.0)"
    ], 11, BLACK)

    # Result
    shape(s, Inches(0.3), Inches(5.2), Inches(9.4), Inches(0.8), NAVY)
    txt(s, Inches(0.5), Inches(5.25), Inches(9), Inches(0.3), "Ms. Sandhya's PDDM Result:", 14, WHITE, True)
    txt(s, Inches(0.5), Inches(5.6), Inches(9), Inches(0.3),
        "Dominant FP = FP1  |  Avg Depth = 2.6 (Mid)  |  Most responses at D3 = strong diagnostic intent",
        12, RGBColor(0xBF, 0xDB, 0xFE))

    # How this feeds SET
    shape(s, Inches(0.3), Inches(6.2), Inches(9.4), Inches(1.0), LIGHT_PURPLE)
    txt(s, Inches(0.5), Inches(6.25), Inches(9), Inches(0.3), "How this feeds SET:", 13, PURPLE, True)
    multi_txt(s, Inches(0.5), Inches(6.6), Inches(9), Inches(0.6), [
        "Her data joins 35 other teachers in branch M031 → school-level orientation distribution",
        "SET1 asks: Does the school's FP distribution match the leaders'? If not, WHERE is the gap?",
        "SET2 asks: Her base001 area scores (A1-A4) explain WHY her depth is at 2.6 and not higher"
    ], 10, BLACK)


def slide_aggregation(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    shape(s, Inches(0), Inches(0), Inches(10), Inches(0.9), TEAL)
    txt(s, Inches(0.5), Inches(0.15), Inches(9), Inches(0.6),
        "From Individual to School Level", 28, WHITE, True)

    txt(s, Inches(0.5), Inches(1.1), Inches(9), Inches(0.3),
        "PDDM gives per-teacher results. SET aggregates them to school/cluster level.", 13, GRAY)

    # Per teacher → Per school
    shape(s, Inches(0.3), Inches(1.7), Inches(4.3), Inches(2.5), LIGHT_BLUE)
    txt(s, Inches(0.5), Inches(1.75), Inches(4), Inches(0.3), "Per Teacher (PDDM)", 14, BLUE, True)
    multi_txt(s, Inches(0.5), Inches(2.15), Inches(4), Inches(2.0), [
        "Teacher A: FP1, Depth 2.6",
        "Teacher B: FP4, Depth 1.8",
        "Teacher C: FP1, Depth 2.1",
        "Teacher D: FP3, Depth 3.0",
        "Teacher E: FP4, Depth 1.5",
        "...(35 teachers total)"
    ], 11, BLACK)

    txt(s, Inches(4.5), Inches(2.7), Inches(1), Inches(0.3), "→", 28, GRAY, True, PP_ALIGN.CENTER)

    shape(s, Inches(5.4), Inches(1.7), Inches(4.3), Inches(2.5), LIGHT_GREEN)
    txt(s, Inches(5.6), Inches(1.75), Inches(4), Inches(0.3), "Per School (SET1)", 14, GREEN, True)
    multi_txt(s, Inches(5.6), Inches(2.15), Inches(4), Inches(2.0), [
        "Teacher FP Distribution:",
        "  FP1: 35%  FP2: 15%  FP3: 20%",
        "  FP4: 25%  FP5: 5%",
        "",
        "Avg Depth by FP:",
        "  FP1: 2.4  FP4: 1.7  FP3: 2.8",
        "",
        "Dominant FP (teachers) = FP1"
    ], 11, BLACK)

    # Plus leader orientation
    shape(s, Inches(0.3), Inches(4.5), Inches(4.3), Inches(1.3), LIGHT_PURPLE)
    txt(s, Inches(0.5), Inches(4.55), Inches(4), Inches(0.3), "Leader Orientation (from 7890/104)", 13, PURPLE, True)
    multi_txt(s, Inches(0.5), Inches(4.95), Inches(4), Inches(0.7), [
        "4 leaders in this branch:",
        "  FP3: 2, FP1: 1, FP4: 1",
        "  Dominant FP (leaders) = FP3"
    ], 11, BLACK)

    txt(s, Inches(4.5), Inches(5.0), Inches(1), Inches(0.3), "→", 28, GRAY, True, PP_ALIGN.CENTER)

    shape(s, Inches(5.4), Inches(4.5), Inches(4.3), Inches(1.3), LIGHT_AMBER)
    txt(s, Inches(5.6), Inches(4.55), Inches(4), Inches(0.3), "Alignment State for FP1:", 13, AMBER, True)
    multi_txt(s, Inches(5.6), Inches(4.95), Inches(4), Inches(0.7), [
        "Teachers: HIGH (35%)",
        "Leaders: LOW (25%)",
        "Depth: LOW (2.4 < 3.0)",
        "→ State: TEACHER-LED INTENT"
    ], 11, BLACK)

    # Connection to base001-004
    shape(s, Inches(0.3), Inches(6.1), Inches(9.4), Inches(1.1), NAVY)
    txt(s, Inches(0.5), Inches(6.15), Inches(9), Inches(0.3), "Then base001-004 explains WHY:", 14, WHITE, True)
    multi_txt(s, Inches(0.5), Inches(6.5), Inches(9), Inches(0.6), [
        "A1 (Learning Diagnostics) = 3.09 (Mid) → tools exist but follow-through is weak",
        "B3 (Teacher Dev Culture) = 3.50 (High) → leaders say growth is protected",
        "A2 (Teacher Dev) = 2.40 (Low) → but teachers don't feel safe → PERCEPTION GAP detected"
    ], 11, RGBColor(0xBF, 0xDB, 0xFE))


def slide_summary(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    shape(s, Inches(0), Inches(0), Inches(10), Inches(7.5), NAVY)
    txt(s, Inches(0.8), Inches(0.8), Inches(8.4), Inches(0.6),
        "PDDM Summary", 32, WHITE, True, PP_ALIGN.CENTER)

    steps = [
        ("1", "SURVEY", "Teacher picks answers (doesn't see tags)", WHITE),
        ("2", "FP TAG", "Each answer has a hidden FP tag (FP1-FP5)", LIGHT_BLUE),
        ("3", "DEPTH TAG", "Each answer has a hidden Depth tag (D1-D4)", LIGHT_GREEN),
        ("4", "DOMINANT FP", "Count FP votes → most frequent = dominant", LIGHT_PURPLE),
        ("5", "AVG DEPTH", "Average depth score → D1-D4 band", LIGHT_AMBER),
        ("6", "PER-FP DEPTH", "Depth for each FP separately (for SET)", LIGHT_TEAL),
        ("7", "AGGREGATE", "All teachers → school-level distribution (SET1)", WHITE),
        ("8", "EXPLAIN", "base001-004 areas explain WHY depth stalls (SET2)", WHITE),
    ]

    y = 1.8
    for num, title, desc, bg in steps:
        shape(s, Inches(1.5), Inches(y), Inches(0.6), Inches(0.45), BLUE)
        txt(s, Inches(1.5), Inches(y + 0.05), Inches(0.6), Inches(0.35), num, 16, WHITE, True, PP_ALIGN.CENTER)
        txt(s, Inches(2.3), Inches(y + 0.02), Inches(1.8), Inches(0.35), title, 14, bg if bg == WHITE else RGBColor(0x93, 0xC5, 0xFD), True)
        txt(s, Inches(4.2), Inches(y + 0.02), Inches(5), Inches(0.35), desc, 12, RGBColor(0xD1, 0xD5, 0xDB))
        y += 0.55

    txt(s, Inches(0.8), Inches(6.3), Inches(8.4), Inches(0.8),
        "PDDM makes invisible intent visible.\nIt turns a simple survey into a diagnostic lens\nwithout teachers ever knowing they're being \"measured.\"",
        15, RGBColor(0xBF, 0xDB, 0xFE), False, PP_ALIGN.CENTER)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_title(prs)
    slide_big_picture(prs)
    slide_fp_explained(prs)
    slide_depth_explained(prs)
    slide_how_tagging_works(prs)
    slide_computation(prs)
    slide_full_example(prs)
    slide_aggregation(prs)
    slide_summary(prs)

    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    prs.save(OUT)
    print(f"Generated: {OUT}")
    print(f"Slides: {len(prs.slides)}")
    print(f"Size: {os.path.getsize(OUT) / 1024:.1f} KB")


if __name__ == "__main__":
    main()
