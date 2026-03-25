#!/usr/bin/env python3
"""
Generate a simple PPTX explaining the SET framework
connected to base001-base004 question sets.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

OUT = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(
    os.path.abspath(__file__)))), "output", "SET_Explained_Simple.pptx")

# Colors
NAVY = RGBColor(0x1E, 0x3A, 0x5F)
BLUE = RGBColor(0x25, 0x63, 0xEB)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
AMBER = RGBColor(0xD9, 0x77, 0x06)
RED = RGBColor(0xDC, 0x26, 0x26)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
TEAL = RGBColor(0x0D, 0x94, 0x88)
GRAY = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF5, 0xF6, 0xFA)
LIGHT_BLUE = RGBColor(0xDB, 0xEA, 0xFE)
LIGHT_GREEN = RGBColor(0xDC, 0xFC, 0xE7)
LIGHT_AMBER = RGBColor(0xFE, 0xF3, 0xC7)
LIGHT_RED = RGBColor(0xFE, 0xE2, 0xE2)
LIGHT_PURPLE = RGBColor(0xED, 0xE9, 0xFE)
BLACK = RGBColor(0x1A, 0x1A, 0x2E)


def add_shape(slide, left, top, width, height, fill_color, text="", font_size=12, font_color=BLACK, bold=False, align=PP_ALIGN.LEFT, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(8)
        tf.margin_right = Pt(8)
        tf.margin_top = Pt(6)
        tf.margin_bottom = Pt(4)
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.alignment = align
    return shape


def add_text(slide, left, top, width, height, text, size=14, color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox


def add_bullet_text(slide, left, top, width, height, lines, size=12, color=BLACK):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(4)
        p.level = 0
    return txBox


def title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_shape(slide, Inches(0), Inches(0), Inches(10), Inches(7.5), NAVY)
    add_text(slide, Inches(0.8), Inches(1.5), Inches(8.4), Inches(1),
             "SET — Structured Explanation of Translation", 36, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, Inches(0.8), Inches(2.5), Inches(8.4), Inches(0.5),
             "A Simple Guide", 24, RGBColor(0xBF, 0xDB, 0xFE), False, PP_ALIGN.CENTER)
    add_text(slide, Inches(0.8), Inches(3.5), Inches(8.4), Inches(1),
             "How NEP intent becomes (or fails to become)\nclassroom practice under real school conditions",
             16, RGBColor(0x93, 0xC5, 0xFD), False, PP_ALIGN.CENTER)
    add_text(slide, Inches(0.8), Inches(5.5), Inches(8.4), Inches(0.5),
             "Connected to: base001 | base002 | base003 | base004",
             14, RGBColor(0x93, 0xC5, 0xFD), False, PP_ALIGN.CENTER)
    add_text(slide, Inches(0.8), Inches(6.5), Inches(8.4), Inches(0.5),
             "Project Kshitij | Myelin | Confidential",
             11, RGBColor(0x60, 0x7D, 0x8B), False, PP_ALIGN.CENTER)


def slide_what_is_set(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, Inches(0), Inches(0), Inches(10), Inches(0.9), NAVY)
    add_text(slide, Inches(0.5), Inches(0.15), Inches(9), Inches(0.6),
             "What is SET?", 28, WHITE, True)

    add_text(slide, Inches(0.5), Inches(1.2), Inches(9), Inches(0.5),
             "SET is NOT a score. It is a structured way to read and explain data.", 16, NAVY, True)

    # Three key questions
    y = 2.0
    for color, q, a in [
        (LIGHT_BLUE, "1. WHERE are the gaps?", "SET1 compares what leaders & teachers intend (orientation) with what actually happens in classrooms (depth)"),
        (LIGHT_GREEN, "2. WHY do gaps exist?", "SET2 looks at school conditions (diagnostic areas) that make practice easier or harder"),
        (LIGHT_PURPLE, "3. WHAT patterns emerge?", "Archetypes — recurring patterns across schools that help target interventions"),
    ]:
        add_shape(slide, Inches(0.5), Inches(y), Inches(9), Inches(1.2), color)
        add_text(slide, Inches(0.7), Inches(y + 0.05), Inches(8.5), Inches(0.4), q, 15, NAVY, True)
        add_text(slide, Inches(0.7), Inches(y + 0.45), Inches(8.5), Inches(0.7), a, 12, GRAY)
        y += 1.4

    add_shape(slide, Inches(0.5), Inches(6.2), Inches(9), Inches(0.8), LIGHT_AMBER)
    add_text(slide, Inches(0.7), Inches(6.3), Inches(8.5), Inches(0.6),
             "Key principle: SET preserves dignity — it explains conditions, not people.\n\"Translation is constrained by specific, addressable conditions — not by teacher resistance.\"",
             11, AMBER, False)


def slide_the_stack(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, Inches(0), Inches(0), Inches(10), Inches(0.9), NAVY)
    add_text(slide, Inches(0.5), Inches(0.15), Inches(9), Inches(0.6),
             "The Interpretive Stack", 28, WHITE, True)

    add_text(slide, Inches(0.5), Inches(1.1), Inches(9), Inches(0.4),
             "Data is read in this order — top to bottom. Each layer explains the next.", 13, GRAY)

    # Stack boxes
    boxes = [
        (BLUE, WHITE, "Layer 1: ORIENTATION", "Where does attention naturally go?", "setCodes: 1234, 7890, 103, 104", "FP1-FP5 (Each Child Unique, Holistic Learning, Reflective Practitioner, Assessment, Collaboration)"),
        (PURPLE, WHITE, "Layer 2: PRACTICE DEPTH", "How far does intent translate into action?", "Derived from orientation responses", "D1 (Procedural) → D2 (Responsive) → D3 (Diagnostic) → D4 (Adaptive)"),
        (TEAL, WHITE, "Layer 3: PRACTICE DIAGNOSTICS", "What conditions make practice easier or harder?", "setCodes: base001, base002, base003, base004", "Teacher: A1-A4  |  Leader: B1-B5"),
    ]

    y = 1.7
    for bg, fg, title, purpose, source, detail in boxes:
        add_shape(slide, Inches(0.5), Inches(y), Inches(9), Inches(1.6), bg)
        add_text(slide, Inches(0.7), Inches(y + 0.05), Inches(4), Inches(0.35), title, 16, fg, True)
        add_text(slide, Inches(0.7), Inches(y + 0.4), Inches(8.5), Inches(0.3), purpose, 13, RGBColor(0xE0, 0xE7, 0xFF))
        add_text(slide, Inches(0.7), Inches(y + 0.75), Inches(8.5), Inches(0.3), "Source: " + source, 10, RGBColor(0xBF, 0xDB, 0xFE))
        add_text(slide, Inches(0.7), Inches(y + 1.05), Inches(8.5), Inches(0.4), detail, 11, RGBColor(0xD1, 0xD5, 0xDB))
        # Arrow
        if y < 5:
            add_text(slide, Inches(4.5), Inches(y + 1.6), Inches(1), Inches(0.3), "▼", 18, GRAY, True, PP_ALIGN.CENTER)
        y += 1.95


def slide_base001_base004(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, Inches(0), Inches(0), Inches(10), Inches(0.9), NAVY)
    add_text(slide, Inches(0.5), Inches(0.15), Inches(9), Inches(0.6),
             "base001 – base004: The Practice Diagnostics", 26, WHITE, True)

    add_text(slide, Inches(0.5), Inches(1.05), Inches(9), Inches(0.4),
             "These are Layer 3 — they explain WHY practice depth looks the way it does.", 13, GRAY)

    # Teacher box
    add_shape(slide, Inches(0.3), Inches(1.6), Inches(4.6), Inches(0.5), BLUE)
    add_text(slide, Inches(0.5), Inches(1.63), Inches(4.2), Inches(0.4),
             "TEACHER LENS  (base001 EN / base003 MR)", 13, WHITE, True)

    teacher_areas = [
        ("A1", "Continuous Learning Diagnostics", "5 Qs", "Can teachers SEE learning gaps?", LIGHT_BLUE),
        ("A2", "Teacher Development & Growth", "5 Qs", "Do teachers feel SAFE to change?", LIGHT_BLUE),
        ("A3", "HPC Enablement", "5 Qs", "Is assessment FORMATIVE or ritual?", LIGHT_BLUE),
        ("A4", "Parent & Community Support", "4 Qs", "Is collaboration ENABLED or dumped?", LIGHT_BLUE),
    ]

    y = 2.2
    for code, name, qs, question, bg in teacher_areas:
        add_shape(slide, Inches(0.3), Inches(y), Inches(4.6), Inches(0.85), bg)
        add_text(slide, Inches(0.5), Inches(y + 0.02), Inches(1), Inches(0.3), code, 14, BLUE, True)
        add_text(slide, Inches(1.2), Inches(y + 0.02), Inches(3.5), Inches(0.3), name + " (" + qs + ")", 11, BLACK, True)
        add_text(slide, Inches(0.5), Inches(y + 0.4), Inches(4.2), Inches(0.4), question, 11, GRAY)
        y += 0.95

    # Leader box
    add_shape(slide, Inches(5.1), Inches(1.6), Inches(4.6), Inches(0.5), PURPLE)
    add_text(slide, Inches(5.3), Inches(1.63), Inches(4.2), Inches(0.4),
             "LEADER LENS  (base002 EN / base004 MR)", 13, WHITE, True)

    leader_areas = [
        ("B1", "NEP Governance & Ownership", "3 Qs", "Is NEP OWNED or just spoken about?", LIGHT_PURPLE),
        ("B2", "Data-Informed Decisions", "3 Qs", "Does data INFLUENCE decisions?", LIGHT_PURPLE),
        ("B3", "Teacher Development Culture", "3 Qs", "Is teacher growth PROTECTED?", LIGHT_PURPLE),
        ("B4", "HPC & Reporting Culture", "3 Qs", "Is HPC a NARRATIVE or compliance?", LIGHT_PURPLE),
        ("B5", "Parent & Community Partnerships", "3 Qs", "Is collaboration STRATEGIC?", LIGHT_PURPLE),
    ]

    y = 2.2
    for code, name, qs, question, bg in leader_areas:
        add_shape(slide, Inches(5.1), Inches(y), Inches(4.6), Inches(0.85), bg)
        add_text(slide, Inches(5.3), Inches(y + 0.02), Inches(1), Inches(0.3), code, 14, PURPLE, True)
        add_text(slide, Inches(6.0), Inches(y + 0.02), Inches(3.5), Inches(0.3), name + " (" + qs + ")", 11, BLACK, True)
        add_text(slide, Inches(5.3), Inches(y + 0.4), Inches(4.2), Inches(0.4), question, 11, GRAY)
        y += 0.95

    # Footer
    add_shape(slide, Inches(0.3), Inches(6.5), Inches(9.4), Inches(0.7), LIGHT_AMBER)
    add_text(slide, Inches(0.5), Inches(6.55), Inches(9), Inches(0.6),
             "19 teacher Qs + 15 leader Qs = 34 scored questions per school\n"
             "All Likert 4-point (Strongly Agree → Strongly Disagree). English + Marathi merged per lens.",
             10, AMBER)


def slide_cross_read(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, Inches(0), Inches(0), Inches(10), Inches(0.9), NAVY)
    add_text(slide, Inches(0.5), Inches(0.15), Inches(9), Inches(0.6),
             "Cross-Reading: Teacher ↔ Leader", 28, WHITE, True)

    add_text(slide, Inches(0.5), Inches(1.1), Inches(9), Inches(0.4),
             "When Teacher and Leader views DIVERGE on the same topic, the gap itself is a signal.", 13, GRAY)

    pairs = [
        ("A2", "B3", "Teacher Development", "Do teachers feel safe to grow?", "Does leadership protect growth?"),
        ("A3", "B4", "HPC / Assessment", "Do teachers understand HPC?", "Does leadership treat HPC as narrative?"),
        ("A4", "B5", "Parent Engagement", "Are teachers supported?", "Does school have a strategy?"),
        ("A1", "B2", "Learning Data", "Can teachers see gaps?", "Does leadership use data for decisions?"),
    ]

    y = 1.7
    for t_area, l_area, topic, t_question, l_question in pairs:
        # Topic label
        add_shape(slide, Inches(3.8), Inches(y), Inches(2.4), Inches(0.4), NAVY)
        add_text(slide, Inches(3.8), Inches(y + 0.02), Inches(2.4), Inches(0.35), topic, 11, WHITE, True, PP_ALIGN.CENTER)
        # Teacher side
        add_shape(slide, Inches(0.3), Inches(y + 0.5), Inches(3.3), Inches(0.8), LIGHT_BLUE)
        add_text(slide, Inches(0.5), Inches(y + 0.5), Inches(1), Inches(0.3), t_area, 14, BLUE, True)
        add_text(slide, Inches(0.5), Inches(y + 0.8), Inches(3), Inches(0.4), t_question, 11, GRAY)
        # Leader side
        add_shape(slide, Inches(6.4), Inches(y + 0.5), Inches(3.3), Inches(0.8), LIGHT_PURPLE)
        add_text(slide, Inches(6.6), Inches(y + 0.5), Inches(1), Inches(0.3), l_area, 14, PURPLE, True)
        add_text(slide, Inches(6.6), Inches(y + 0.8), Inches(3), Inches(0.4), l_question, 11, GRAY)
        y += 1.4

    # Divergence key
    add_shape(slide, Inches(0.3), Inches(y + 0.1), Inches(9.4), Inches(0.9), LIGHT_AMBER)
    add_text(slide, Inches(0.5), Inches(y + 0.15), Inches(9), Inches(0.8),
             "Both High = Validated strength  |  Both Low = Confirmed gap\n"
             "Leader High + Teacher Low = Perception gap (leaders think it's fine, teachers don't)\n"
             "Teacher High + Leader Low = Informal support exists despite no formal structure",
             10, AMBER)


def slide_set1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, Inches(0), Inches(0), Inches(10), Inches(0.9), BLUE)
    add_text(slide, Inches(0.5), Inches(0.15), Inches(9), Inches(0.6),
             "SET1: WHERE Are the Gaps?", 28, WHITE, True)

    add_text(slide, Inches(0.5), Inches(1.1), Inches(9), Inches(0.4),
             "SET1 compares orientation (what people intend) with depth (what actually happens).", 13, GRAY)

    # FP table
    add_text(slide, Inches(0.5), Inches(1.7), Inches(3), Inches(0.3), "5 Foundational Philosophies (FPs)", 13, NAVY, True)
    fps = [
        ("FP1", "Each Child is Unique"),
        ("FP2", "Holistic & Experiential Learning"),
        ("FP3", "Teacher as Reflective Practitioner"),
        ("FP4", "Assessment for Learning"),
        ("FP5", "Collaboration & Community"),
    ]
    y = 2.1
    for code, name in fps:
        add_shape(slide, Inches(0.5), Inches(y), Inches(4), Inches(0.35), LIGHT_BLUE)
        add_text(slide, Inches(0.6), Inches(y), Inches(1), Inches(0.3), code, 11, BLUE, True)
        add_text(slide, Inches(1.3), Inches(y), Inches(3), Inches(0.3), name, 11, BLACK)
        y += 0.4

    # Depth table
    add_text(slide, Inches(5.1), Inches(1.7), Inches(3), Inches(0.3), "4 Practice Depth Levels", 13, NAVY, True)
    depths = [
        ("D1", "Procedural", "Follows steps"),
        ("D2", "Responsive", "Adjusts to students"),
        ("D3", "Diagnostic", "Identifies root causes"),
        ("D4", "Adaptive", "Creates new approaches"),
    ]
    colors = [LIGHT_RED, LIGHT_AMBER, LIGHT_BLUE, LIGHT_GREEN]
    y = 2.1
    for (code, name, desc), col in zip(depths, colors):
        add_shape(slide, Inches(5.1), Inches(y), Inches(4.5), Inches(0.35), col)
        add_text(slide, Inches(5.2), Inches(y), Inches(0.6), Inches(0.3), code, 11, NAVY, True)
        add_text(slide, Inches(5.8), Inches(y), Inches(1.5), Inches(0.3), name, 11, BLACK, True)
        add_text(slide, Inches(7.3), Inches(y), Inches(2.2), Inches(0.3), desc, 10, GRAY)
        y += 0.4

    # Alignment states
    add_text(slide, Inches(0.5), Inches(4.3), Inches(9), Inches(0.3), "5 Alignment States (per FP, per school)", 13, NAVY, True)

    states = [
        ("Intent High, Translation Low", "Both intend it, but classrooms haven't changed", LIGHT_RED, RED),
        ("Teacher-Led Intent", "Teachers ahead of leadership", LIGHT_AMBER, AMBER),
        ("Leader-Led Intent", "Leadership ahead of teachers", LIGHT_PURPLE, PURPLE),
        ("Low Attention Zone", "Neither side focuses here", LIGHT_RED, RED),
        ("Aligned", "Intent + depth + practice all match", LIGHT_GREEN, GREEN),
    ]

    y = 4.7
    for name, desc, bg, fg in states:
        add_shape(slide, Inches(0.5), Inches(y), Inches(9), Inches(0.4), bg)
        add_text(slide, Inches(0.7), Inches(y + 0.02), Inches(3.5), Inches(0.3), name, 11, fg, True)
        add_text(slide, Inches(4.2), Inches(y + 0.02), Inches(5), Inches(0.3), desc, 11, GRAY)
        y += 0.45


def slide_set2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, Inches(0), Inches(0), Inches(10), Inches(0.9), TEAL)
    add_text(slide, Inches(0.5), Inches(0.15), Inches(9), Inches(0.6),
             "SET2: WHY Do Gaps Exist?", 28, WHITE, True)

    add_text(slide, Inches(0.5), Inches(1.1), Inches(9), Inches(0.5),
             "SET2 uses the diagnostic areas from base001-base004 to explain\nwhy practice depth stalls at a given level.", 13, GRAY)

    # FP → Area links
    add_text(slide, Inches(0.5), Inches(1.9), Inches(9), Inches(0.3),
             "Each FP has specific areas that explain its depth:", 13, NAVY, True)

    links = [
        ("FP1", "Each Child Unique", "A1 + A3", "B2 + B4", "Can teachers see uniqueness? Is HPC alive?"),
        ("FP2", "Holistic Learning", "A3 + A2", "B1 + B4", "Shared meaning? Safe to experiment?"),
        ("FP3", "Reflective Practitioner", "A2", "B3", "Is reflection safe and supported?"),
        ("FP4", "Assessment for Learning", "A1 + A3", "B2 + B4", "Tools exist? Data used?"),
        ("FP5", "Collaboration", "A4", "B1 + B3", "Is collaboration enabled? Ownership shared?"),
    ]

    y = 2.4
    # Header
    for x, text, w in [(0.5, "FP", 0.5), (1.1, "Philosophy", 2.2), (3.4, "Teacher Areas", 1.5), (5.0, "Leader Areas", 1.5), (6.6, "Core Question", 3.2)]:
        add_shape(slide, Inches(x), Inches(y), Inches(w), Inches(0.35), NAVY)
        add_text(slide, Inches(x + 0.05), Inches(y), Inches(w), Inches(0.3), text, 10, WHITE, True)
    y += 0.4

    for fp, name, t_areas, l_areas, question in links:
        bg = LIGHT_BG if links.index((fp, name, t_areas, l_areas, question)) % 2 == 0 else WHITE
        for x, text, w in [(0.5, fp, 0.5), (1.1, name, 2.2), (3.4, t_areas, 1.5), (5.0, l_areas, 1.5), (6.6, question, 3.2)]:
            add_shape(slide, Inches(x), Inches(y), Inches(w), Inches(0.4), bg)
            add_text(slide, Inches(x + 0.05), Inches(y + 0.02), Inches(w - 0.1), Inches(0.35), text, 10, BLACK, x == 0.5)
        y += 0.45

    # Banding
    add_text(slide, Inches(0.5), Inches(5.1), Inches(9), Inches(0.3),
             "Area Score Banding (mean of sub-area Likert scores, 1-4):", 13, NAVY, True)

    for x, label, score_range, bg, fg in [
        (0.5, "LOW", "< 2.5", LIGHT_RED, RED),
        (3.5, "MID", "2.5 – 3.2", LIGHT_AMBER, AMBER),
        (6.5, "HIGH", "> 3.2", LIGHT_GREEN, GREEN),
    ]:
        add_shape(slide, Inches(x), Inches(5.5), Inches(2.8), Inches(0.7), bg)
        add_text(slide, Inches(x + 0.1), Inches(5.5), Inches(2.6), Inches(0.3), label + " (" + score_range + ")", 13, fg, True, PP_ALIGN.CENTER)
        descs = {"LOW": "System not supporting", "MID": "Partial support, gaps exist", "HIGH": "Actively supported"}
        add_text(slide, Inches(x + 0.1), Inches(5.85), Inches(2.6), Inches(0.3), descs[label], 10, GRAY, False, PP_ALIGN.CENTER)

    add_shape(slide, Inches(0.5), Inches(6.5), Inches(9), Inches(0.6), LIGHT_AMBER)
    add_text(slide, Inches(0.7), Inches(6.55), Inches(8.6), Inches(0.5),
             "SET2 Rule: Areas explain depth — they do NOT define orientation.\n"
             "Teacher areas (A) = classroom constraints. Leader areas (B) = system constraints. Never merge at individual level.",
             10, AMBER)


def slide_example(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, Inches(0), Inches(0), Inches(10), Inches(0.9), NAVY)
    add_text(slide, Inches(0.5), Inches(0.15), Inches(9), Inches(0.6),
             "Example: Reading a Branch", 28, WHITE, True)

    add_text(slide, Inches(0.5), Inches(1.1), Inches(9), Inches(0.3),
             "Branch M017 — how SET reads the data step by step:", 14, NAVY, True)

    steps = [
        ("1", "ORIENTATION", "Teachers' dominant FP = FP4 (Assessment for Learning)\nLeaders' dominant FP = FP3 (Reflective Practitioner)", LIGHT_BLUE, BLUE),
        ("2", "DEPTH", "Teacher practice depth for FP4 = D2 (Responsive)\nDepth hasn't reached D3 (Diagnostic) yet", LIGHT_AMBER, AMBER),
        ("3", "DIAGNOSTICS\n(base001)", "A1 = 3.09 (Mid) — tools exist but not fully used\nA3 = 2.97 (Mid) — HPC understanding partial\nA1 Follow-through = LOW — no time to act on insights", LIGHT_RED, RED),
        ("4", "SET SENTENCE", "\"Given high FP4 orientation and practice depth at D2,\ntranslation is shaped by low follow-through in A1,\nindicating that time constraints — not teacher\nresistance — are the limiting factor.\"", LIGHT_GREEN, GREEN),
    ]

    y = 1.6
    for num, title, content, bg, fg in steps:
        add_shape(slide, Inches(0.5), Inches(y), Inches(0.6), Inches(0.5), fg)
        add_text(slide, Inches(0.5), Inches(y + 0.05), Inches(0.6), Inches(0.4), num, 18, WHITE, True, PP_ALIGN.CENTER)
        add_shape(slide, Inches(1.2), Inches(y), Inches(8.3), Inches(1.2 if num == "4" else 0.95), bg)
        add_text(slide, Inches(1.3), Inches(y + 0.02), Inches(2), Inches(0.3), title, 12, fg, True)
        add_text(slide, Inches(1.3), Inches(y + 0.3), Inches(8), Inches(0.8), content, 11, BLACK)
        y += 1.35 if num == "4" else 1.1


def slide_archetypes(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, Inches(0), Inches(0), Inches(10), Inches(0.9), PURPLE)
    add_text(slide, Inches(0.5), Inches(0.15), Inches(9), Inches(0.6),
             "6 Archetypes: Common Patterns", 28, WHITE, True)

    add_text(slide, Inches(0.5), Inches(1.05), Inches(9), Inches(0.3),
             "Archetypes help schools recognize their pattern and target interventions.", 13, GRAY)

    archetypes = [
        ("ARCH-01", "Intent without tools", "Everyone intends it, but no shared language or tools", LIGHT_RED),
        ("ARCH-02", "Reflection without safety", "Teachers reflect but don't act — fear blocks change", LIGHT_RED),
        ("ARCH-03", "Teachers ahead of systems", "Teachers lead, but leadership hasn't operationalised support", LIGHT_AMBER),
        ("ARCH-04", "Leadership ahead of teachers", "Leadership intends it, but classrooms don't have follow-through", LIGHT_AMBER),
        ("ARCH-05", "Low attention zone", "Neither side focuses here — foundational orientation needed", LIGHT_RED),
        ("ARCH-06", "Fully aligned", "Intent + depth + conditions all coherent — sustain and deepen", LIGHT_GREEN),
    ]

    y = 1.5
    for arch_id, label, desc, bg in archetypes:
        add_shape(slide, Inches(0.5), Inches(y), Inches(9), Inches(0.75), bg)
        add_text(slide, Inches(0.7), Inches(y + 0.02), Inches(1.3), Inches(0.3), arch_id, 12, PURPLE, True)
        add_text(slide, Inches(2.1), Inches(y + 0.02), Inches(3), Inches(0.3), label, 12, BLACK, True)
        add_text(slide, Inches(2.1), Inches(y + 0.35), Inches(7), Inches(0.3), desc, 11, GRAY)
        y += 0.85

    add_shape(slide, Inches(0.5), Inches(6.6), Inches(9), Inches(0.6), LIGHT_PURPLE)
    add_text(slide, Inches(0.7), Inches(6.65), Inches(8.6), Inches(0.5),
             "Archetypes are sense-making constructs, NOT evaluations.\n"
             "They are group-level only (school/cluster). They never imply ranking or blame.",
             10, PURPLE)


def slide_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, Inches(0), Inches(0), Inches(10), Inches(7.5), NAVY)
    add_text(slide, Inches(0.8), Inches(0.8), Inches(8.4), Inches(0.6),
             "Summary: The Complete Picture", 30, WHITE, True, PP_ALIGN.CENTER)

    flow = [
        "Orientation Surveys (1234/7890/103/104)",
        "    ↓  reveal FP1-FP5 intent",
        "Practice Depth (derived from orientation)",
        "    ↓  shows D1-D4 translation level",
        "Practice Diagnostics (base001-base004)",
        "    ↓  explain conditions via A1-A4 + B1-B5",
        "SET1: WHERE gaps exist",
        "    ↓  alignment states per FP",
        "SET2: WHY gaps exist",
        "    ↓  area conditions + canonical sentences",
        "Archetypes: WHAT pattern this is",
        "    →  targeted intervention design",
    ]

    y = 1.8
    for line in flow:
        is_arrow = line.strip().startswith("↓") or line.strip().startswith("→")
        add_text(slide, Inches(1.5), Inches(y), Inches(7), Inches(0.35),
                 line, 13 if not is_arrow else 11,
                 RGBColor(0x93, 0xC5, 0xFD) if is_arrow else WHITE,
                 not is_arrow, PP_ALIGN.LEFT)
        y += 0.35

    add_text(slide, Inches(0.8), Inches(6.2), Inches(8.4), Inches(0.8),
             "SET tells schools: \"Here's what's happening, here's why,\nand here's what would make the biggest difference.\"",
             15, RGBColor(0xBF, 0xDB, 0xFE), False, PP_ALIGN.CENTER)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    title_slide(prs)
    slide_what_is_set(prs)
    slide_the_stack(prs)
    slide_base001_base004(prs)
    slide_cross_read(prs)
    slide_set1(prs)
    slide_set2(prs)
    slide_example(prs)
    slide_archetypes(prs)
    slide_summary(prs)

    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    prs.save(OUT)
    print(f"Generated: {OUT}")
    print(f"Slides: {len(prs.slides)}")
    print(f"Size: {os.path.getsize(OUT) / 1024:.1f} KB")


if __name__ == "__main__":
    main()
