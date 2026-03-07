#!/usr/bin/env python3
"""Generate SRI Methodology documents — detailed formula explanations in plain language.

Outputs:
  docs/SRI_Methodology_Explained.html   — styled for Cmd+P → Save as PDF
  docs/SRI_Methodology_Explained.pptx   — ~18-slide 16:9 presentation
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

BASE = Path(__file__).resolve().parent.parent
DOCS = BASE / "docs"
DOCS.mkdir(exist_ok=True)

TITLE = "How We Measure NEP Readiness: The School Readiness Index (SRI)"
SUBTITLE = "Deccan Education Society — Kshitij Pilot"
DATE = "February 2026"


# ═══════════════════════════════════════════════════════════════
#  Part 1 — HTML (styled for print-to-PDF)
# ═══════════════════════════════════════════════════════════════

def _esc(text):
    return text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


def generate_html():
    html = f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<title>{TITLE}</title>
<style>
  @page {{ size: A4; margin: 20mm 18mm 20mm 18mm; }}
  body {{
    font-family: -apple-system, 'Helvetica Neue', Arial, sans-serif;
    font-size: 12.5px; line-height: 1.7; color: #1e293b;
    max-width: 720px; margin: 0 auto; padding: 30px 20px;
  }}
  h1 {{
    font-size: 22px; color: #1e3a5f; text-align: center;
    border-bottom: 3px solid #2563eb; padding-bottom: 10px; margin-bottom: 4px;
  }}
  .subtitle {{ text-align: center; font-size: 12px; color: #6b7280; margin-bottom: 28px; }}
  h2 {{
    font-size: 16px; color: #1e40af; margin-top: 30px; padding-bottom: 4px;
    border-bottom: 1px solid #cbd5e1; page-break-after: avoid;
  }}
  h3 {{
    font-size: 13.5px; color: #334155; margin-top: 18px; margin-bottom: 4px;
    page-break-after: avoid;
  }}
  .maxtag {{
    font-size: 10px; font-weight: 400; color: #6b7280;
    background: #f1f5f9; padding: 2px 8px; border-radius: 10px;
    margin-left: 6px; vertical-align: middle;
  }}
  p {{ margin: 6px 0; }}
  ul {{ padding-left: 20px; margin: 6px 0; }}
  li {{ margin-bottom: 4px; }}
  .analogy {{
    background: #eff6ff; border: 1px solid #bfdbfe; border-radius: 6px;
    padding: 10px 14px; margin: 10px 0; font-style: italic; color: #1e40af;
    font-size: 12px;
  }}
  .step-box {{
    background: #f8fafc; border-left: 3px solid #2563eb; padding: 10px 14px;
    margin: 8px 0; font-size: 12px; line-height: 1.65;
  }}
  .step-box.purple {{ border-left-color: #8b5cf6; }}
  .step-box.green {{ border-left-color: #10b981; }}
  .step-box.amber {{ border-left-color: #f59e0b; }}
  .step-box.gray {{ border-left-color: #6b7280; }}
  .formula-plain {{
    background: #fefce8; border: 1px solid #fde68a; border-radius: 6px;
    padding: 10px 14px; margin: 10px 0; font-size: 12px;
    font-family: 'SF Mono', Consolas, monospace; line-height: 1.8;
  }}
  .example-box {{
    background: #f0fdf4; border: 1px solid #bbf7d0; border-radius: 6px;
    padding: 10px 14px; margin: 10px 0; font-size: 12px;
  }}
  .example-box strong {{ color: #166534; }}
  .note-box {{
    background: #fff7ed; border: 1px solid #fed7aa; border-radius: 6px;
    padding: 8px 14px; margin: 10px 0; font-size: 11.5px; color: #9a3412;
  }}
  table.simple {{
    border-collapse: collapse; margin: 10px 0; font-size: 11.5px;
    width: 100%;
  }}
  table.simple th {{
    background: #f1f5f9; font-weight: 600; text-align: left;
    padding: 6px 10px; border: 1px solid #e2e8f0;
  }}
  table.simple td {{
    padding: 5px 10px; border: 1px solid #e2e8f0;
  }}
  .footer {{
    margin-top: 40px; padding-top: 12px; border-top: 1px solid #e2e8f0;
    font-size: 10px; color: #9ca3af; text-align: center;
  }}
</style>
</head>
<body>
<h1>{_esc(TITLE)}</h1>
<div class="subtitle">{_esc(SUBTITLE)} &mdash; {_esc(DATE)}</div>

<!-- ═══ SECTION: What is SRI? ═══ -->
<h2>What is SRI?</h2>
<div class="analogy">Think of a full-body health checkup: the doctor checks your heart, lungs, blood, bones, and reflexes separately, then gives you one overall health picture. SRI does the same for a school's NEP readiness.</div>
<p>The <strong>School Readiness Index</strong> gives each school branch a single score from <strong>0 to 100</strong> that captures how ready it is for the National Education Policy (NEP 2020). The score is built from 5 independent areas — each looking at a different dimension of readiness.</p>

<!-- ═══ SECTION: 5 Areas ═══ -->
<h2>The 5 Areas at a Glance</h2>
<table class="simple">
<tr><th>Area</th><th>What It Measures</th><th>Max Points</th></tr>
<tr><td><strong>C1 Intent</strong></td><td>Do teachers and leaders share the same NEP vision?</td><td>20</td></tr>
<tr><td><strong>C2 Practice</strong></td><td>How deeply are teachers applying NEP ideas in class?</td><td>25</td></tr>
<tr><td><strong>C3 Capacity</strong></td><td>Can teachers think and solve problems at higher levels?</td><td>20</td></tr>
<tr><td><strong>C4 System</strong></td><td>Does the school's culture and processes support NEP?</td><td>20</td></tr>
<tr><td><strong>C5 Ecosystem</strong></td><td>Are parents and community involved? <em>(not yet measured)</em></td><td>15</td></tr>
<tr><td colspan="2" style="text-align:right"><strong>Total</strong></td><td><strong>100</strong></td></tr>
</table>
<p>Because C5 is not yet measured, the current maximum any branch can achieve is <strong>85 out of 100</strong>.</p>

<!-- ═══ SECTION: C1 INTENT ═══ -->
<h2>C1 — Intent: Do teachers and leaders share the same NEP vision? <span class="maxtag">max 20 pts</span></h2>
<div class="analogy">Analogy: A rowing team. Even the strongest rowers won't win if they're pulling in different directions. This score checks whether everyone is rowing together.</div>

<h3>What data do we use?</h3>
<p>We surveyed both <strong>teachers</strong> and <strong>school leaders</strong> about five NEP focus areas (called "Focus Points"):</p>
<ul>
<li><strong>FP1</strong> — Each Child is Unique</li>
<li><strong>FP2</strong> — Competency-Based Learning</li>
<li><strong>FP3</strong> — Teacher Upskilling</li>
<li><strong>FP4</strong> — Diagnosing Learning Levels</li>
<li><strong>FP5</strong> — Parent-Teacher Collaboration</li>
</ul>
<p>Each survey question belongs to one of these 5 focus points. Respondents pick an option on a 0–3 scale (0 = lowest, 3 = highest aspiration).</p>

<h3>How we calculate it — step by step</h3>

<div class="step-box">
<strong>Step 1 — Depth Score per focus area</strong><br>
For each focus point, we average all survey responses and divide by 3 (the maximum). This gives a score from 0 to 1 that tells us "how deeply do people feel about this area?"<br><br>
<em>Example:</em> If 10 teachers answered FP1 questions with responses averaging 2.4 out of 3, the depth score = 2.4 / 3 = <strong>0.80</strong>. This means teachers have 80% of the maximum possible aspiration for "Each Child is Unique".
</div>

<div class="step-box">
<strong>Step 2 — Coverage: did people address all 5 areas?</strong><br>
We count how many of the 5 focus points have a meaningful depth score (above 10%). Then divide by 5.<br><br>
<em>Example:</em> If all 5 focus points have depth scores above 0.10, coverage = 5/5 = <strong>1.00</strong> (100%). If only 3 areas were addressed, coverage = 3/5 = <strong>0.60</strong> (60%).
</div>

<div class="step-box">
<strong>Step 3 — Balance: are responses spread evenly?</strong><br>
We check whether the 5 depth scores are roughly equal or whether one area dominates and others are neglected. We compare the spread (variance) against the worst possible spread.<br><br>
<em>Example:</em> Scores of [0.80, 0.75, 0.82, 0.78, 0.85] are very balanced — the number comes out close to <strong>1.00</strong>. Scores of [0.95, 0.10, 0.05, 0.90, 0.10] are very lopsided — balance drops to near <strong>0.00</strong>.
</div>

<div class="step-box">
<strong>Step 4 — Alignment: do teachers and leaders agree?</strong><br>
We compare the teacher's pattern of 5 scores with the leader's pattern. Think of each as an arrow pointing in some direction. If both arrows point the same way (same priority order), alignment is near <strong>1.0</strong>. If they point in very different directions, alignment drops toward <strong>0.0</strong>.<br><br>
Technically, this uses "cosine similarity" — but all that means in plain language is: <em>do both groups rank the 5 areas in roughly the same order of importance?</em>
</div>

<div class="step-box">
<strong>Step 5 — Combine into the final C1 score</strong><br>
For teachers: multiply Coverage &times; Balance &times; Alignment = Teacher Score (0 to 1).<br>
For leaders: multiply Coverage &times; Balance &times; Alignment = Leader Score (0 to 1).<br><br>
<strong>C1 = 10 &times; Teacher Score + 10 &times; Leader Score</strong><br><br>
Teachers contribute up to 10 points, leaders contribute up to 10 points. The maximum is 20.
</div>

<div class="formula-plain">
Depth Score per FP = average response / 3<br>
Coverage = (FPs above threshold) / 5<br>
Balance = 1 &minus; (how spread-out the scores are) / (worst possible spread)<br>
Alignment = how similarly teachers and leaders rank the 5 areas<br><br>
Teacher Score = Coverage &times; Balance &times; Alignment<br>
Leader Score = Coverage &times; Balance &times; Alignment<br>
<strong>C1 = 10 &times; Teacher Score + 10 &times; Leader Score</strong>
</div>

<div class="example-box">
<strong>Worked example — Branch M001 (Ahilyadevi High School):</strong><br>
Teacher score = 0.989 (excellent coverage, balance, and alignment)<br>
Leader score = 0.971 (similarly strong)<br>
Alignment between them = 0.994 (nearly identical priorities)<br>
C1 = 10 &times; 0.989 + 10 &times; 0.971 = <strong>19.6 out of 20</strong>
</div>

<div class="note-box">
<strong>Important:</strong> If a branch has only teachers or only leaders (but not both), the alignment score is 0. This forces C1 to 0 — by design. Intent readiness means <em>both</em> roles must be present and aligned.
</div>

<!-- ═══ SECTION: C2 PRACTICE ═══ -->
<h2>C2 — Practice: How deeply are teachers practising NEP ideas? <span class="maxtag">max 25 pts</span></h2>
<div class="analogy">Analogy: Learning to swim. First you wade in the shallow end (Level 1). Then you float (Level 2). Then you swim laps (Level 3). Then you dive in the deep end (Level 4). Deeper practice = higher score.</div>

<h3>What data do we use?</h3>
<p>The same teacher survey as C1, but now we interpret each answer as a <strong>depth level</strong>:</p>
<table class="simple">
<tr><th>Response</th><th>Depth</th><th>Meaning</th></tr>
<tr><td>Option 0</td><td>Level 1</td><td>"I'm aware of it"</td></tr>
<tr><td>Option 1</td><td>Level 2</td><td>"I'm starting to try it"</td></tr>
<tr><td>Option 2</td><td>Level 3</td><td>"I do it regularly"</td></tr>
<tr><td>Option 3</td><td>Level 4</td><td>"It shapes my teaching deeply"</td></tr>
</table>

<h3>How we calculate it — step by step</h3>

<div class="step-box purple">
<strong>Step 1 — Build a depth grid</strong><br>
For each focus point (FP1–FP5) and each depth level (1–4), count how many teachers are at that level. This creates a 5-by-4 grid.
</div>

<div class="step-box purple">
<strong>Step 2 — Remove impossible combinations</strong><br>
Some focus-point + depth-level combinations are impossible to observe in a real classroom. For example, you can't realistically see "surface-level awareness" of "Each Child is Unique" in actual practice — it's either visible or it isn't. These cells are zeroed out:<br>
&bull; FP1 at Level 1 (too vague to observe)<br>
&bull; FP1 at Level 4 (not achievable in current context)<br>
&bull; FP2 at Level 4 (same reason)
</div>

<div class="step-box purple">
<strong>Step 3 — Average depth per focus area</strong><br>
For each focus point, calculate the average depth level of teachers (using only the observable cells). Higher average = teachers are practising at deeper levels.<br><br>
<em>Example:</em> If for FP3, 5 teachers are at Level 2, 10 at Level 3, and 5 at Level 4, the average depth = (5&times;2 + 10&times;3 + 5&times;4) / 20 = <strong>3.0</strong>.
</div>

<div class="step-box purple">
<strong>Step 4 — Normalise to a 0-to-1 scale</strong><br>
Each focus area has a different maximum achievable depth (because of the removed cells). FP1 and FP2 can only go up to Level 3 effectively (max = 2.0); FP3–FP5 can reach Level 4 (max = 3.0). We divide each area's average by its maximum, giving a 0-to-1 score.
</div>

<div class="step-box purple">
<strong>Step 5 — Combine into the final C2 score</strong><br>
Coverage = how many of the 5 focus areas have data, divided by 5.<br>
Depth Effectiveness = Coverage &times; average of the normalised depth scores.<br><br>
<strong>C2 = 25 &times; Depth Effectiveness</strong>
</div>

<div class="formula-plain">
Average depth per FP = weighted average of (level &times; count) / total<br>
Normalised depth = average depth / max achievable depth for that FP<br>
Coverage = (FPs with data) / 5<br>
Depth Effectiveness = Coverage &times; average(normalised depths)<br>
<strong>C2 = 25 &times; Depth Effectiveness</strong>
</div>

<div class="example-box">
<strong>Worked example — Branch M001:</strong><br>
All 5 FPs covered. Depth Effectiveness = 0.546 (teachers are roughly at the mid-range of depth).<br>
C2 = 25 &times; 0.546 = <strong>13.65 out of 25</strong><br>
<em>Interpretation: Teachers show moderate depth. There is room to move from "I know about it" to "I actively do it."</em>
</div>

<!-- ═══ SECTION: C3 CAPACITY ═══ -->
<h2>C3 — Capacity: Can teachers think at higher levels? <span class="maxtag">max 20 pts</span></h2>
<div class="analogy">Analogy: A cognitive fitness test. Just as a gym tests your strength, speed, and flexibility, MathTangle tests how well teachers handle challenging thinking — not whether they know maths content, but how they think.</div>

<h3>What data do we use?</h3>
<p><strong>MathTangle</strong> — a short adaptive test of maths puzzles. The test automatically gets harder or easier based on each teacher's answers. 107 teachers across 28 branches took the test.</p>

<h3>Four thinking dimensions</h3>
<p>Each teacher gets a score from 0 to 1 on four dimensions, then the four are averaged:</p>

<div class="step-box green">
<strong>Dimension 1 — Higher-Order Thinking</strong><br>
Of all the questions the teacher answered, what fraction required <em>analysis</em> or <em>evaluation</em> (rather than just remembering or applying a formula)?<br><br>
<em>In plain terms:</em> We count how many "hard thinking" questions they got right, out of the total. More higher-order correct answers = higher score.
</div>

<div class="step-box green">
<strong>Dimension 2 — Tackling Harder Problems</strong><br>
The test has three difficulty tiers: Level 1 (easy), Level 2 (medium), Level 3 (hard). We average the teacher's accuracy on Level 2 and Level 3.<br><br>
<em>In plain terms:</em> Can the teacher handle the tough questions? If they get 70% right on medium and 50% right on hard, the score = (0.70 + 0.50) / 2 = <strong>0.60</strong>.
</div>

<div class="step-box green">
<strong>Dimension 3 — Thoughtful Revision</strong><br>
The test tracks an "indecision score" — how much a teacher hesitated, went back and forth, or second-guessed themselves. We compare each teacher's indecision to the <em>most indecisive teacher across all branches</em>.<br><br>
<em>In plain terms:</em> Score = 1 minus (your indecision / worst indecision). Low hesitation = high score. This rewards <em>deliberate, confident</em> thinking.
</div>

<div class="step-box green">
<strong>Dimension 4 — Consistency</strong><br>
Did the teacher frequently change their answer back and forth ("flip-flops")? We count flip-flops relative to total questions.<br><br>
<em>In plain terms:</em> Score = 1 minus (flip-flops / total questions). If a teacher answered 30 questions and flipped 3, the score = 1 &minus; 3/30 = <strong>0.90</strong>. Fewer flip-flops = more consistent, reliable thinking.
</div>

<div class="step-box green">
<strong>Combining the 4 dimensions</strong><br>
Each teacher's Cognitive Proficiency Index (CPI) = simple average of the 4 dimension scores.<br>
Branch CPI = average of all teachers' CPIs in that branch.<br><br>
<strong>C3 = 20 &times; Branch CPI</strong>
</div>

<div class="formula-plain">
Dim 1 = (Analyse correct + Evaluate correct) / total questions<br>
Dim 2 = (Level-2 accuracy + Level-3 accuracy) / 2<br>
Dim 3 = 1 &minus; (your indecision / max indecision across all)<br>
Dim 4 = 1 &minus; (flip-flops / total questions)<br><br>
CPI per teacher = average of Dim 1, 2, 3, 4<br>
CPI per branch = average of all teachers' CPIs<br>
<strong>C3 = 20 &times; Branch CPI</strong>
</div>

<div class="example-box">
<strong>Worked example — Branch M001 (7 teachers tested):</strong><br>
Average CPI across 7 teachers = 0.763<br>
C3 = 20 &times; 0.763 = <strong>15.26 out of 20</strong><br>
<em>Interpretation: Teachers show solid cognitive capacity with room for improvement on higher-order and harder problems.</em>
</div>

<!-- ═══ SECTION: C4 SYSTEM ═══ -->
<h2>C4 — System: Does the school culture support NEP? <span class="maxtag">max 20 pts</span></h2>
<div class="analogy">Analogy: A plant can only grow if the soil is right. Similarly, teachers can only implement NEP well if the school's systems, culture, and leadership support them.</div>

<h3>What data do we use?</h3>
<p>A survey where teachers and leaders respond to statements with <strong>Strongly Agree / Agree / Disagree / Strongly Disagree</strong>. 676 respondents, 12,576 individual response rows.</p>

<h3>How we calculate it — step by step</h3>

<div class="step-box amber">
<strong>Step 1 — Convert to Yes/No</strong><br>
Each response is simplified: <strong>Strongly Agree or Agree = 1</strong> (positive). <strong>Disagree or Strongly Disagree = 0</strong> (negative).<br><br>
<em>Why simplify?</em> A 4-point scale can be noisy. The key question is: does this person fundamentally agree or disagree?
</div>

<div class="step-box amber">
<strong>Step 2 — Group questions by area</strong><br>
Teacher questions fall into 4 areas:<br>
&bull; <strong>A1</strong> — Continuous Learning Diagnostics<br>
&bull; <strong>A2</strong> — Teacher Development<br>
&bull; <strong>A3</strong> — Holistic Progress Card<br>
&bull; <strong>A4</strong> — Parent &amp; Community<br><br>
Leader questions fall into 5 areas:<br>
&bull; <strong>B1</strong> — NEP Governance<br>
&bull; <strong>B2</strong> — Data-Informed Decisions<br>
&bull; <strong>B3</strong> — Teacher Development Culture<br>
&bull; <strong>B4</strong> — Reporting &amp; Progress Cards<br>
&bull; <strong>B5</strong> — Parent &amp; Community Partnerships
</div>

<div class="step-box amber">
<strong>Step 3 — Average within each area</strong><br>
For each area, calculate the percentage of positive responses. For example, if 90% of A1 responses are Agree/Strongly Agree, then A1 score = <strong>0.90</strong>.
</div>

<div class="step-box amber">
<strong>Step 4 — Combine areas into two indices</strong><br>
<strong>Enablement Index</strong> = average of A1, A2, A3, A4 scores. This represents "how well do teachers feel supported?"<br>
<strong>Routine Index</strong> = average of B1, B2, B3, B4, B5 scores. This represents "how strong are the leadership routines?"
</div>

<div class="step-box amber">
<strong>Step 5 — Final score</strong><br>
<strong>C4 = 10 &times; Enablement Index + 10 &times; Routine Index</strong><br><br>
Teacher perspective contributes up to 10 points, leader perspective up to 10 points. Maximum is 20.
</div>

<div class="formula-plain">
Each response: Agree/SA &rarr; 1, Disagree/SD &rarr; 0<br>
Area mean = sum(binary) / count for each area (A1..A4, B1..B5)<br>
Enablement Index = average(A1, A2, A3, A4)<br>
Routine Index = average(B1, B2, B3, B4, B5)<br>
<strong>C4 = 10 &times; Enablement + 10 &times; Routine</strong>
</div>

<div class="example-box">
<strong>Worked example — Branch M001:</strong><br>
Enablement Index (teachers) = 0.911 (91% positive responses across A1-A4)<br>
Routine Index (leaders) = 0.850 (85% positive across B1-B5)<br>
C4 = 10 &times; 0.911 + 10 &times; 0.850 = <strong>17.61 out of 20</strong>
</div>

<div class="note-box">
<strong>Important:</strong> Branches without leader respondents get Routine Index = 0, capping C4 at roughly 10 out of 20. Both roles must participate for a full system score.
</div>

<!-- ═══ SECTION: C5 ECOSYSTEM ═══ -->
<h2>C5 — Ecosystem: Are parents and community involved? <span class="maxtag">max 15 pts</span></h2>
<div class="analogy">Education doesn't stop at the school gate — the whole village matters.</div>
<p>This area will eventually measure parent, student, and community engagement with NEP. Planned instruments:</p>
<ul>
<li><strong>Parent Perception Survey</strong> — How do parents view the school's NEP efforts?</li>
<li><strong>Student Voice Survey</strong> — Do students notice changes in how they learn?</li>
<li><strong>Community Partnership Audit</strong> — Is the wider community engaged?</li>
</ul>
<p>Since these are <strong>not yet deployed</strong>, every branch currently scores <strong>0 out of 15</strong>.</p>

<!-- ═══ SECTION: Assembly ═══ -->
<h2>Putting It All Together</h2>
<div class="formula-plain">
<strong>SRI = C1 + C2 + C3 + C4 + C5</strong><br>
Maximum (with all 5): 20 + 25 + 20 + 20 + 15 = <strong>100</strong><br>
Current maximum (C5 = 0): 20 + 25 + 20 + 20 + 0 = <strong>85</strong>
</div>

<h3>What do the scores mean?</h3>
<table class="simple">
<tr><th>Score Range</th><th>Band</th><th>What It Means</th></tr>
<tr><td>60–85 / 85 (70%+)</td><td><strong>Strong</strong></td><td>Well on its way — continue and refine</td></tr>
<tr><td>43–59 / 85 (50–70%)</td><td><strong>Developing</strong></td><td>Good foundation but clear gaps to address</td></tr>
<tr><td>&lt; 43 / 85 (&lt;50%)</td><td><strong>Needs Attention</strong></td><td>Focused effort needed in multiple areas</td></tr>
</table>

<div class="example-box">
<strong>Full worked example — Branch M001 (Ahilyadevi High School):</strong><br><br>
C1 (Intent) = 19.60 / 20<br>
C2 (Practice) = 13.65 / 25<br>
C3 (Capacity) = 15.26 / 20<br>
C4 (System) = 17.61 / 20<br>
C5 (Ecosystem) = 0.00 / 15<br>
<strong>SRI Total = 66.12 / 85 (78%)</strong> &mdash; Strong readiness
</div>

<!-- ═══ SECTION: Data Coverage ═══ -->
<h2>Data at a Glance</h2>
<table class="simple">
<tr><th>Construct</th><th>Data Source</th><th>Respondents</th></tr>
<tr><td>C1 Intent</td><td>Teacher + leader intent surveys (EN &amp; MR)</td><td>~717 responses, 29 branches with non-zero C1</td></tr>
<tr><td>C2 Practice</td><td>Same teacher intent surveys (depth analysis)</td><td>36 branches scored</td></tr>
<tr><td>C3 Capacity</td><td>MathTangle adaptive cognitive test</td><td>107 teachers, 28 branches</td></tr>
<tr><td>C4 System</td><td>Agree/disagree system survey</td><td>676 respondents, 12,576 response rows, 36 branches</td></tr>
<tr><td>C5 Ecosystem</td><td>Not yet deployed</td><td>—</td></tr>
</table>
<p><strong>37 DES branches</strong> scored in total. <strong>25 branches</strong> have non-zero scores in all 4 measured constructs. Average SRI among those 25 is approximately <strong>65 / 85 (~76%)</strong>, putting them solidly in the Strong readiness band.</p>
<p><em>Excluded:</em> M008 (demo branch) and m942e (internal test school).</p>

<div class="footer">Myelin NEP Readiness Pilot &mdash; {DATE} &mdash; For internal use</div>
</body>
</html>"""

    path = DOCS / "SRI_Methodology_Explained.html"
    path.write_text(html, encoding="utf-8")
    print(f"HTML -> {path}")
    return path


# ═══════════════════════════════════════════════════════════════
#  Part 2 — PPTX (16:9, ~18 slides)
# ═══════════════════════════════════════════════════════════════

PRIMARY   = RGBColor(0x1A, 0x36, 0x5D)
ACCENT    = RGBColor(0x2B, 0x6C, 0xB0)
GREEN     = RGBColor(0x38, 0xA1, 0x69)
AMBER     = RGBColor(0xD6, 0x9E, 0x2E)
RED       = RGBColor(0xE5, 0x3E, 0x3E)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xF7, 0xFA, 0xFC)
MUTED     = RGBColor(0x71, 0x80, 0x96)
TEXT_CLR  = RGBColor(0x2D, 0x37, 0x48)
C1_CLR    = RGBColor(0x3B, 0x82, 0xF6)
C2_CLR    = RGBColor(0x8B, 0x5C, 0xF6)
C3_CLR    = RGBColor(0x10, 0xB9, 0x81)
C4_CLR    = RGBColor(0xF5, 0x9E, 0x0B)
C5_CLR    = RGBColor(0x6B, 0x72, 0x80)
SRI_CLR   = RGBColor(0x1E, 0x3A, 0x5F)
INFO_BG   = RGBColor(0xEF, 0xF6, 0xFF)
INFO_CLR  = RGBColor(0x1E, 0x40, 0xAF)
EX_BG     = RGBColor(0xF0, 0xFD, 0xF4)
EX_CLR    = RGBColor(0x16, 0x65, 0x34)
WARN_BG   = RGBColor(0xFF, 0xF7, 0xED)
WARN_CLR  = RGBColor(0x9A, 0x34, 0x12)


def add_bg(slide, color=LIGHT_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=WHITE, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def set_text(shape, text, size=12, bold=False, color=TEXT_CLR, alignment=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return tf


def add_para(tf, text, size=12, bold=False, color=TEXT_CLR, space_before=4):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.space_before = Pt(space_before)
    return p


def title_bar(slide, text, color=PRIMARY):
    s = add_shape(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8), color)
    set_text(s, f"  {text}", size=24, bold=True, color=WHITE)


def step_card(slide, left, top, width, height, num, title, body, accent_clr):
    """Numbered step card with accent stripe."""
    bx = add_shape(slide, left, top, width, height)
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Pt(5), height)
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent_clr
    stripe.line.fill.background()
    tf = set_text(bx, f"    Step {num}: {title}", size=13, bold=True, color=PRIMARY)
    add_para(tf, f"    {body}", size=11, color=TEXT_CLR, space_before=6)
    return bx


def generate_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── Slide 1: Title ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PRIMARY)
    box = add_shape(slide, Inches(1), Inches(1.2), Inches(11.3), Inches(5), PRIMARY)
    box.line.fill.background()
    tf = set_text(box, TITLE, size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_para(tf, "", size=14, space_before=16)
    add_para(tf, "A detailed guide — formulas explained in everyday language", size=20, color=RGBColor(0x93, 0xC5, 0xFD), space_before=12)
    add_para(tf, "", size=14, space_before=20)
    add_para(tf, SUBTITLE, size=16, color=RGBColor(0x93, 0xC5, 0xFD), space_before=8)
    add_para(tf, DATE, size=14, color=RGBColor(0x93, 0xC5, 0xFD), space_before=8)

    # ── Slide 2: What is SRI? + 5 Areas ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_bar(slide, "What is SRI?")

    box = add_shape(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(1.8))
    tf = set_text(box, "  The School Readiness Index is like a health checkup for your school.", size=18, bold=True, color=PRIMARY)
    add_para(tf, "  Instead of checking blood pressure, it checks how ready your school is for NEP 2020.", size=14, color=TEXT_CLR, space_before=8)
    add_para(tf, "  Each branch gets a single score from 0 to 100, built from 5 independent areas.", size=14, color=ACCENT, space_before=8)

    # Analogy box
    box = add_shape(slide, Inches(0.5), Inches(3.5), Inches(12.3), Inches(0.6), INFO_BG, C1_CLR)
    set_text(box, '  Like a fitness test checks strength, flexibility, stamina separately, then combines into one picture.', size=12, color=INFO_CLR)

    # 5-area KPI boxes
    constructs = [
        ("C1 Intent", "20", C1_CLR), ("C2 Practice", "25", C2_CLR),
        ("C3 Capacity", "20", C3_CLR), ("C4 System", "20", C4_CLR),
        ("C5 Ecosystem", "15", C5_CLR),
    ]
    for i, (name, pts, clr) in enumerate(constructs):
        left = Inches(0.5 + i * 2.55)
        bx = add_shape(slide, left, Inches(4.4), Inches(2.2), Inches(1.6))
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(4.4), Inches(2.2), Pt(5))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = clr
        stripe.line.fill.background()
        tf = set_text(bx, pts, size=36, bold=True, color=clr, alignment=PP_ALIGN.CENTER)
        add_para(tf, name, size=12, color=MUTED, space_before=6)
        tf.paragraphs[-1].alignment = PP_ALIGN.CENTER

    bx = add_shape(slide, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.4))
    set_text(bx, "  Total = 100 pts   |   Currently achievable: 85 (Ecosystem not yet measured)   |   SRI = C1 + C2 + C3 + C4 + C5", size=12, bold=True, color=SRI_CLR, alignment=PP_ALIGN.CENTER)

    # ── Slide 3: C1 Intent Overview ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_bar(slide, "C1 — Intent: Do teachers and leaders share the same NEP vision?", C1_CLR)

    box = add_shape(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(1.3))
    tf = set_text(box, "  We surveyed both teachers and leaders about 5 NEP Focus Points:", size=14, color=TEXT_CLR)
    add_para(tf, "  FP1: Each Child is Unique  |  FP2: Competency-Based Learning  |  FP3: Teacher Upskilling", size=13, color=ACCENT, space_before=8)
    add_para(tf, "  FP4: Diagnosing Learning Levels  |  FP5: Parent-Teacher Collaboration", size=13, color=ACCENT, space_before=4)

    box = add_shape(slide, Inches(0.5), Inches(3.0), Inches(12.3), Inches(0.6), INFO_BG, C1_CLR)
    set_text(box, '  Analogy: A rowing team — even the strongest rowers won\'t win if they\'re pulling in different directions.', size=12, color=INFO_CLR)

    checks = [
        ("Coverage", "How many of the 5 areas\ndid people respond about?", "If all 5: score = 1.0\nIf only 3: score = 0.6", C1_CLR),
        ("Balance", "Are responses spread evenly\nacross the 5 areas?", "Evenly spread: score ~1.0\nOne area dominates: ~0.0", ACCENT),
        ("Alignment", "Do teachers and leaders\nagree on priorities?", "Same priority order: ~1.0\nVery different: ~0.0", GREEN),
    ]
    for i, (label, desc, ex, clr) in enumerate(checks):
        left = Inches(0.5 + i * 4.2)
        bx = add_shape(slide, left, Inches(3.9), Inches(3.8), Inches(2.8))
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(3.9), Inches(3.8), Pt(5))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = clr
        stripe.line.fill.background()
        tf = set_text(bx, label, size=18, bold=True, color=clr, alignment=PP_ALIGN.CENTER)
        add_para(tf, desc, size=12, color=TEXT_CLR, space_before=12)
        tf.paragraphs[-1].alignment = PP_ALIGN.CENTER
        add_para(tf, ex, size=11, color=MUTED, space_before=12)
        tf.paragraphs[-1].alignment = PP_ALIGN.CENTER

    box = add_shape(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4))
    set_text(box, "  C1 = 10 x Teacher(Coverage x Balance x Alignment) + 10 x Leader(Coverage x Balance x Alignment)  |  Max = 20", size=11, bold=True, color=PRIMARY, alignment=PP_ALIGN.CENTER)

    # ── Slide 4: C1 How the Math Works ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_bar(slide, "C1 — How the Math Works (in plain language)", C1_CLR)

    step_card(slide, Inches(0.5), Inches(1.4), Inches(5.8), Inches(1.3), "1", "Depth Score per area",
              "For each focus point, average all responses (0-3) and divide\nby 3. Gives a 0-to-1 score. E.g., average response 2.4/3 = 0.80.", C1_CLR)
    step_card(slide, Inches(6.8), Inches(1.4), Inches(5.8), Inches(1.3), "2", "Coverage",
              "Count how many of the 5 areas got a score above 10%,\nthen divide by 5. All 5 active = 5/5 = 1.0 (100%).", C1_CLR)
    step_card(slide, Inches(0.5), Inches(2.9), Inches(5.8), Inches(1.3), "3", "Balance",
              "Compare how spread-out the 5 scores are vs. the worst\npossible spread. Equal scores = ~1.0. Lopsided = ~0.0.", C1_CLR)
    step_card(slide, Inches(6.8), Inches(2.9), Inches(5.8), Inches(1.3), "4", "Alignment",
              "Compare teacher priorities vs. leader priorities.\n\"Do both groups rank the 5 areas in the same order?\"", C1_CLR)
    step_card(slide, Inches(0.5), Inches(4.4), Inches(12.3), Inches(0.9), "5", "Combine",
              "Teacher Score = Coverage x Balance x Alignment.  Leader Score = same.  C1 = 10 x Teacher + 10 x Leader.", C1_CLR)

    # Example box
    box = add_shape(slide, Inches(0.5), Inches(5.6), Inches(12.3), Inches(1.3), EX_BG, GREEN)
    tf = set_text(box, "  Worked Example — Branch M001 (Ahilyadevi High School)", size=14, bold=True, color=EX_CLR)
    add_para(tf, "  Teacher score = 0.989 (excellent coverage, balance, alignment)", size=12, color=TEXT_CLR, space_before=6)
    add_para(tf, "  Leader score = 0.971  |  Alignment between them = 0.994 (nearly identical priorities)", size=12, color=TEXT_CLR, space_before=4)
    add_para(tf, "  C1 = 10 x 0.989 + 10 x 0.971 = 19.6 out of 20", size=13, bold=True, color=EX_CLR, space_before=6)

    # Note
    box = add_shape(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4), WARN_BG)
    set_text(box, "  If a branch has only teachers OR only leaders (not both), Alignment = 0 and C1 = 0. Both roles must be present.", size=11, color=WARN_CLR)

    # ── Slide 5: C2 Practice Overview ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_bar(slide, "C2 — Practice: How deeply are teachers practising NEP ideas?", C2_CLR)

    box = add_shape(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.8))
    tf = set_text(box, "  Same teacher survey as C1, but now each answer is read as a depth level.", size=14, color=TEXT_CLR)
    add_para(tf, "  Higher depth = teachers are not just aware, they're actively doing it in class.", size=14, color=ACCENT, space_before=6)

    box = add_shape(slide, Inches(0.5), Inches(2.5), Inches(12.3), Inches(0.5), RGBColor(0xF5, 0xF3, 0xFF), C2_CLR)
    set_text(box, '  Analogy: Learning to swim — first you wade (Level 1), then float (2), then swim laps (3), then dive deep (4).', size=12, color=RGBColor(0x55, 0x3C, 0x9A))

    depths = [
        ("Level 1", "Surface awareness", "\"I've heard of it\"", RGBColor(0xBF, 0xDB, 0xFE)),
        ("Level 2", "Beginning practice", "\"I'm starting to try\"", RGBColor(0x93, 0xC5, 0xFD)),
        ("Level 3", "Active practice", "\"I do it regularly\"", RGBColor(0x60, 0xA5, 0xFA)),
        ("Level 4", "Deep integration", "\"It shapes my teaching\"", C2_CLR),
    ]
    for i, (lvl, label, quote, clr) in enumerate(depths):
        left = Inches(0.5 + i * 3.2)
        bx = add_shape(slide, left, Inches(3.3), Inches(2.8), Inches(1.5))
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(3.3), Inches(2.8), Pt(5))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = clr
        stripe.line.fill.background()
        tf = set_text(bx, lvl, size=16, bold=True, color=PRIMARY, alignment=PP_ALIGN.CENTER)
        add_para(tf, label, size=12, color=TEXT_CLR, space_before=8)
        tf.paragraphs[-1].alignment = PP_ALIGN.CENTER
        add_para(tf, quote, size=11, color=MUTED, space_before=6)
        tf.paragraphs[-1].alignment = PP_ALIGN.CENTER

    box = add_shape(slide, Inches(0.5), Inches(5.1), Inches(12.3), Inches(0.8))
    tf = set_text(box, "  Some focus-area + level combinations can't be observed in real classrooms:", size=12, color=MUTED)
    add_para(tf, "  FP1 at Level 1 (too vague to see)  |  FP1 at Level 4 (not achievable yet)  |  FP2 at Level 4 (same)", size=12, color=MUTED, space_before=4)

    box = add_shape(slide, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.5))
    set_text(box, "  C2 = 25 x (fraction of FPs covered) x (average normalised depth across FPs)  |  Max = 25", size=12, bold=True, color=PRIMARY, alignment=PP_ALIGN.CENTER)

    # ── Slide 6: C2 How the Math Works ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_bar(slide, "C2 — How the Math Works (in plain language)", C2_CLR)

    step_card(slide, Inches(0.5), Inches(1.4), Inches(5.8), Inches(1.3), "1", "Build a depth grid",
              "For each FP (1-5) and depth (1-4), count how many\nteachers are at that level. Creates a 5 x 4 grid.", C2_CLR)
    step_card(slide, Inches(6.8), Inches(1.4), Inches(5.8), Inches(1.3), "2", "Remove impossible cells",
              "Zero out combinations that can't be observed in class:\nFP1-Level1, FP1-Level4, FP2-Level4.", C2_CLR)
    step_card(slide, Inches(0.5), Inches(2.9), Inches(5.8), Inches(1.3), "3", "Average depth per FP",
              "Weighted average of levels (using observable cells only).\nHigher average = teachers practise at deeper levels.", C2_CLR)
    step_card(slide, Inches(6.8), Inches(2.9), Inches(5.8), Inches(1.3), "4", "Normalise to 0-1",
              "Each FP has different max depth (FP1/2: max 2.0,\nFP3-5: max 3.0). Scale each to a 0-to-1 range.", C2_CLR)
    step_card(slide, Inches(0.5), Inches(4.4), Inches(12.3), Inches(0.9), "5", "Combine into Depth Effectiveness",
              "Coverage = (FPs with data) / 5.  Depth Effectiveness = Coverage x average(normalised depths).  C2 = 25 x Depth Effectiveness.", C2_CLR)

    box = add_shape(slide, Inches(0.5), Inches(5.6), Inches(12.3), Inches(1.0), EX_BG, GREEN)
    tf = set_text(box, "  Worked Example — Branch M001", size=14, bold=True, color=EX_CLR)
    add_para(tf, "  All 5 FPs covered. Depth Effectiveness = 0.546 (mid-range — teachers aware but not deeply practising yet).", size=12, color=TEXT_CLR, space_before=6)
    add_para(tf, "  C2 = 25 x 0.546 = 13.65 out of 25", size=13, bold=True, color=EX_CLR, space_before=6)

    # ── Slide 7: C3 Capacity Overview ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_bar(slide, "C3 — Capacity: Can teachers think at higher levels?", C3_CLR)

    box = add_shape(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.8))
    tf = set_text(box, "  MathTangle — an adaptive maths test that gets harder or easier based on each teacher's answers.", size=14, color=TEXT_CLR)
    add_para(tf, "  107 teachers across 28 branches. Tests thinking ability, not maths content knowledge.", size=13, color=ACCENT, space_before=6)

    box = add_shape(slide, Inches(0.5), Inches(2.5), Inches(12.3), Inches(0.5), RGBColor(0xEC, 0xFD, 0xF5), C3_CLR)
    set_text(box, '  Analogy: A cognitive fitness test — measuring thinking strength, not textbook knowledge.', size=12, color=RGBColor(0x06, 0x5F, 0x46))

    dims = [
        ("Higher-Order\nThinking", "What fraction of answers\nwere at the harder 'analyse'\nand 'evaluate' levels?", C3_CLR),
        ("Tackling Harder\nProblems", "Average accuracy on\nthe medium and hard\ndifficulty questions.", ACCENT),
        ("Thoughtful\nRevision", "Low indecision = good.\nRewards deliberate,\nconfident thinking.", AMBER),
        ("Consistency", "Few flip-flops = good.\nSteady performance\nthroughout the test.", RGBColor(0x63, 0x66, 0xF1)),
    ]
    for i, (label, desc, clr) in enumerate(dims):
        left = Inches(0.5 + i * 3.2)
        bx = add_shape(slide, left, Inches(3.3), Inches(2.8), Inches(2.6))
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(3.3), Inches(2.8), Pt(5))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = clr
        stripe.line.fill.background()
        tf = set_text(bx, label, size=14, bold=True, color=PRIMARY, alignment=PP_ALIGN.CENTER)
        add_para(tf, desc, size=11, color=TEXT_CLR, space_before=10)
        tf.paragraphs[-1].alignment = PP_ALIGN.CENTER

    box = add_shape(slide, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.5))
    set_text(box, "  Each dimension scores 0-1. CPI per teacher = average of all 4. C3 = 20 x branch average CPI  |  Max = 20", size=12, bold=True, color=PRIMARY, alignment=PP_ALIGN.CENTER)

    # ── Slide 8: C3 How the Math Works ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_bar(slide, "C3 — How the Math Works (in plain language)", C3_CLR)

    step_card(slide, Inches(0.5), Inches(1.4), Inches(5.8), Inches(1.4), "1", "Higher-Order Thinking",
              "Count 'Analyse' + 'Evaluate' questions answered correctly.\nDivide by total cognitive questions.\nMore hard-thinking correct = higher score.", C3_CLR)
    step_card(slide, Inches(6.8), Inches(1.4), Inches(5.8), Inches(1.4), "2", "Tackling Harder Problems",
              "Test has 3 difficulty tiers (Easy, Medium, Hard).\nAverage accuracy on Medium and Hard.\nE.g., 70% Medium + 50% Hard = (0.7+0.5)/2 = 0.60.", C3_CLR)
    step_card(slide, Inches(0.5), Inches(3.0), Inches(5.8), Inches(1.4), "3", "Thoughtful Revision",
              "Score = 1 minus (your indecision / worst indecision).\nCompared to the most indecisive teacher across ALL branches.\nLow hesitation = high score = confident thinking.", C3_CLR)
    step_card(slide, Inches(6.8), Inches(3.0), Inches(5.8), Inches(1.4), "4", "Consistency",
              "Score = 1 minus (flip-flops / total questions).\n30 questions, 3 flip-flops = 1 - 3/30 = 0.90.\nFewer random changes = more reliable thinking.", C3_CLR)

    step_card(slide, Inches(0.5), Inches(4.6), Inches(12.3), Inches(0.8), "5", "Combine",
              "CPI per teacher = average of 4 dimensions.  Branch CPI = average across all teachers.  C3 = 20 x Branch CPI.", C3_CLR)

    box = add_shape(slide, Inches(0.5), Inches(5.7), Inches(12.3), Inches(1.0), EX_BG, GREEN)
    tf = set_text(box, "  Worked Example — Branch M001 (7 teachers tested)", size=14, bold=True, color=EX_CLR)
    add_para(tf, "  Average CPI across 7 teachers = 0.763 (solid but room for improvement on higher-order questions).", size=12, color=TEXT_CLR, space_before=6)
    add_para(tf, "  C3 = 20 x 0.763 = 15.26 out of 20", size=13, bold=True, color=EX_CLR, space_before=6)

    # ── Slide 9: C4 System Overview ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_bar(slide, "C4 — System: Does the school culture support NEP?", C4_CLR)

    box = add_shape(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.8))
    tf = set_text(box, "  Survey: teachers and leaders respond Strongly Agree / Agree / Disagree / Strongly Disagree.", size=14, color=TEXT_CLR)
    add_para(tf, "  676 respondents, 12,576 individual response rows.", size=13, color=ACCENT, space_before=6)

    box = add_shape(slide, Inches(0.5), Inches(2.5), Inches(12.3), Inches(0.5), RGBColor(0xFF, 0xFB, 0xEB), C4_CLR)
    set_text(box, '  Analogy: Checking if the soil is right for the plant — teachers can only thrive if the school system supports them.', size=12, color=RGBColor(0x92, 0x40, 0x0E))

    t_areas = ["A1 — Continuous Learning Diagnostics", "A2 — Teacher Development", "A3 — Holistic Progress Card", "A4 — Parent & Community"]
    l_areas = ["B1 — NEP Governance", "B2 — Data-Informed Decisions", "B3 — Teacher Development Culture", "B4 — Reporting & Progress Cards", "B5 — Parent & Community Partnerships"]

    bx = add_shape(slide, Inches(0.5), Inches(3.3), Inches(5.8), Inches(3.0))
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3.3), Pt(5), Inches(3.0))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = C4_CLR
    stripe.line.fill.background()
    tf = set_text(bx, "  Teacher Areas (Enablement)", size=15, bold=True, color=PRIMARY)
    for a in t_areas:
        add_para(tf, f"    {a}", size=12, color=TEXT_CLR, space_before=8)
    add_para(tf, "    Enablement Index = average of A1, A2, A3, A4", size=11, bold=True, color=C4_CLR, space_before=12)

    bx = add_shape(slide, Inches(6.8), Inches(3.3), Inches(6.0), Inches(3.0))
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(3.3), Pt(5), Inches(3.0))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = ACCENT
    stripe.line.fill.background()
    tf = set_text(bx, "  Leader Areas (Routine)", size=15, bold=True, color=PRIMARY)
    for a in l_areas:
        add_para(tf, f"    {a}", size=12, color=TEXT_CLR, space_before=6)
    add_para(tf, "    Routine Index = average of B1, B2, B3, B4, B5", size=11, bold=True, color=ACCENT, space_before=8)

    box = add_shape(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.5))
    set_text(box, "  C4 = 10 x Enablement Index + 10 x Routine Index  |  Teachers up to 10 pts, Leaders up to 10 pts  |  Max = 20", size=12, bold=True, color=PRIMARY, alignment=PP_ALIGN.CENTER)

    # ── Slide 10: C4 How the Math Works ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_bar(slide, "C4 — How the Math Works (in plain language)", C4_CLR)

    step_card(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(1.1), "1", "Convert every answer to Yes/No",
              "Strongly Agree or Agree = 1 (positive).  Disagree or Strongly Disagree = 0 (negative).\nWhy? A 4-point scale can be noisy. The key question: do they fundamentally agree or not?", C4_CLR)
    step_card(slide, Inches(0.5), Inches(2.7), Inches(12.3), Inches(1.1), "2", "Average within each area",
              "Group questions by area (A1, A2, ... B5). Calculate % positive.\nE.g., if 90% of A1 responses are Agree/SA, then A1 score = 0.90.", C4_CLR)
    step_card(slide, Inches(0.5), Inches(4.0), Inches(5.8), Inches(1.1), "3", "Enablement Index",
              "Average of A1, A2, A3, A4 scores.\n\"How well do teachers feel supported?\"", C4_CLR)
    step_card(slide, Inches(6.8), Inches(4.0), Inches(5.8), Inches(1.1), "4", "Routine Index",
              "Average of B1, B2, B3, B4, B5 scores.\n\"How strong are leadership routines?\"", C4_CLR)

    box = add_shape(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.2), EX_BG, GREEN)
    tf = set_text(box, "  Worked Example — Branch M001", size=14, bold=True, color=EX_CLR)
    add_para(tf, "  Enablement Index (teachers) = 0.911 (91% positive across A1-A4)", size=12, color=TEXT_CLR, space_before=6)
    add_para(tf, "  Routine Index (leaders) = 0.850 (85% positive across B1-B5)", size=12, color=TEXT_CLR, space_before=4)
    add_para(tf, "  C4 = 10 x 0.911 + 10 x 0.850 = 17.61 out of 20", size=13, bold=True, color=EX_CLR, space_before=6)

    box = add_shape(slide, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4), WARN_BG)
    set_text(box, "  Branches without leader respondents get Routine Index = 0, capping C4 at roughly 10/20. Both roles must participate.", size=11, color=WARN_CLR)

    # ── Slide 11: C5 Ecosystem ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_bar(slide, "C5 — Ecosystem: Are parents and community involved?", C5_CLR)

    box = add_shape(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.8))
    tf = set_text(box, "  Education doesn't stop at the school gate — the whole village matters.", size=16, bold=True, color=PRIMARY)
    add_para(tf, "  This construct will measure engagement from parents, students, and the wider community.", size=14, color=TEXT_CLR, space_before=6)

    instruments = [
        ("Parent Perception\nSurvey", "How do parents view\nthe school's NEP efforts?", RGBColor(0x63, 0x66, 0xF1)),
        ("Student Voice\nSurvey", "Do students feel changes\nin how they learn?", GREEN),
        ("Community Partnership\nAudit", "Is the wider community\nengaged with the school?", AMBER),
    ]
    for i, (label, desc, clr) in enumerate(instruments):
        left = Inches(0.5 + i * 4.2)
        bx = add_shape(slide, left, Inches(2.6), Inches(3.8), Inches(2.0))
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(2.6), Inches(3.8), Pt(5))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = clr
        stripe.line.fill.background()
        tf = set_text(bx, label, size=15, bold=True, color=PRIMARY, alignment=PP_ALIGN.CENTER)
        add_para(tf, desc, size=12, color=MUTED, space_before=12)
        tf.paragraphs[-1].alignment = PP_ALIGN.CENTER

    box = add_shape(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.8), RGBColor(0xF3, 0xF4, 0xF6))
    tf = set_text(box, "  Status: Not yet deployed. All branches score 0 / 15.", size=15, bold=True, color=C5_CLR, alignment=PP_ALIGN.CENTER)
    add_para(tf, "  Once these instruments are in place, the maximum SRI rises from 85 to 100.", size=12, color=MUTED, space_before=6)
    tf.paragraphs[-1].alignment = PP_ALIGN.CENTER

    # ── Slide 12: Putting It All Together ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_bar(slide, "Putting It All Together")

    box = add_shape(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(1.3))
    tf = set_text(box, "  SRI = C1 + C2 + C3 + C4 + C5", size=22, bold=True, color=PRIMARY, alignment=PP_ALIGN.CENTER)
    add_para(tf, "  (up to 20) + (up to 25) + (up to 20) + (up to 20) + (up to 15) = up to 100", size=14, color=ACCENT, space_before=8)
    tf.paragraphs[-1].alignment = PP_ALIGN.CENTER
    add_para(tf, "  With C5 not yet measured, current max = 85", size=13, color=MUTED, space_before=6)
    tf.paragraphs[-1].alignment = PP_ALIGN.CENTER

    bands = [
        ("70%+ of max", "60-85 / 85", "Strong", "Well on its way", GREEN),
        ("50-70%", "43-59 / 85", "Developing", "Good base, gaps to address", AMBER),
        ("Below 50%", "< 43 / 85", "Needs Attention", "Focused effort needed", RED),
    ]
    for i, (pct, rng, label, desc, clr) in enumerate(bands):
        left = Inches(0.5 + i * 4.2)
        bx = add_shape(slide, left, Inches(3.1), Inches(3.8), Inches(2.0))
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(3.1), Inches(3.8), Pt(5))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = clr
        stripe.line.fill.background()
        tf = set_text(bx, pct, size=20, bold=True, color=clr, alignment=PP_ALIGN.CENTER)
        add_para(tf, rng, size=12, color=MUTED, space_before=6)
        tf.paragraphs[-1].alignment = PP_ALIGN.CENTER
        add_para(tf, label, size=15, bold=True, color=PRIMARY, space_before=10)
        tf.paragraphs[-1].alignment = PP_ALIGN.CENTER
        add_para(tf, desc, size=11, color=MUTED, space_before=4)
        tf.paragraphs[-1].alignment = PP_ALIGN.CENTER

    box = add_shape(slide, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.2), EX_BG, GREEN)
    tf = set_text(box, "  Full Worked Example — Branch M001 (Ahilyadevi High School for Girls)", size=14, bold=True, color=EX_CLR)
    add_para(tf, "  C1 = 19.60/20   C2 = 13.65/25   C3 = 15.26/20   C4 = 17.61/20   C5 = 0/15", size=13, color=TEXT_CLR, space_before=6)
    add_para(tf, "  SRI Total = 66.12 / 85 (78%) — Strong readiness band", size=14, bold=True, color=EX_CLR, space_before=6)

    # ── Slide 13: Data at a Glance ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_bar(slide, "Data at a Glance")

    stats = [
        ("~717", "Intent Survey\nResponses", "Teachers + Leaders,\nEN & MR", C1_CLR),
        ("107", "MathTangle\nTest Takers", "Adaptive cognitive\nassessment", C3_CLR),
        ("676", "System Survey\nRespondents", "12,576 individual\nresponse rows", C4_CLR),
    ]
    for i, (val, label, desc, clr) in enumerate(stats):
        left = Inches(0.5 + i * 4.2)
        bx = add_shape(slide, left, Inches(1.5), Inches(3.8), Inches(2.3))
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(1.5), Inches(3.8), Pt(5))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = clr
        stripe.line.fill.background()
        tf = set_text(bx, val, size=40, bold=True, color=clr, alignment=PP_ALIGN.CENTER)
        add_para(tf, label, size=13, bold=True, color=PRIMARY, space_before=8)
        tf.paragraphs[-1].alignment = PP_ALIGN.CENTER
        add_para(tf, desc, size=11, color=MUTED, space_before=6)
        tf.paragraphs[-1].alignment = PP_ALIGN.CENTER

    branch_stats = [
        ("37", "DES Branches\nScored", ACCENT),
        ("25", "Branches with\nAll 4 Constructs", GREEN),
        ("2", "Excluded\n(demo/test)", RED),
    ]
    for i, (val, label, clr) in enumerate(branch_stats):
        left = Inches(0.5 + i * 4.2)
        bx = add_shape(slide, left, Inches(4.2), Inches(3.8), Inches(2.0))
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(4.2), Inches(3.8), Pt(5))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = clr
        stripe.line.fill.background()
        tf = set_text(bx, val, size=36, bold=True, color=clr, alignment=PP_ALIGN.CENTER)
        add_para(tf, label, size=12, color=MUTED, space_before=6)
        tf.paragraphs[-1].alignment = PP_ALIGN.CENTER

    box = add_shape(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.5))
    set_text(box, "  Average SRI among 25 complete branches: ~65/85 (76%) — Strong readiness band", size=13, bold=True, color=SRI_CLR, alignment=PP_ALIGN.CENTER)

    # ── Slide 14: How the 5 Areas Connect ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_bar(slide, "How the 5 Areas Connect")

    connections = [
        ("C1 Intent", "What do we\nbelieve in?", "Survey of\nteachers + leaders", C1_CLR),
        ("C2 Practice", "How deeply are\nwe doing it?", "Same survey,\ndepth analysis", C2_CLR),
        ("C3 Capacity", "Can we think\nat higher levels?", "MathTangle\nadaptive test", C3_CLR),
        ("C4 System", "Does the system\nsupport us?", "Agree/disagree\nsurvey", C4_CLR),
        ("C5 Ecosystem", "Is everyone\ninvolved?", "Coming soon", C5_CLR),
    ]
    for i, (name, question, source, clr) in enumerate(connections):
        left = Inches(0.5 + i * 2.55)
        bx = add_shape(slide, left, Inches(1.5), Inches(2.2), Inches(4.5))
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(1.5), Inches(2.2), Pt(5))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = clr
        stripe.line.fill.background()
        tf = set_text(bx, name, size=15, bold=True, color=clr, alignment=PP_ALIGN.CENTER)
        add_para(tf, question, size=13, color=PRIMARY, space_before=14)
        tf.paragraphs[-1].alignment = PP_ALIGN.CENTER
        add_para(tf, "", size=8, space_before=10)
        add_para(tf, "Data source:", size=10, bold=True, color=MUTED, space_before=6)
        tf.paragraphs[-1].alignment = PP_ALIGN.CENTER
        add_para(tf, source, size=11, color=MUTED, space_before=4)
        tf.paragraphs[-1].alignment = PP_ALIGN.CENTER

    box = add_shape(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.5))
    set_text(box, "  Intent  +  Practice  +  Capacity  +  System  +  Ecosystem  =  SRI Total", size=16, bold=True, color=SRI_CLR, alignment=PP_ALIGN.CENTER)

    # ── Slide 15: Formula Summary (cheat sheet) ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_bar(slide, "Formula Summary — All Constructs at a Glance")

    formulas = [
        ("C1 Intent", "20", C1_CLR,
         "Depth Score = avg response / 3 for each FP\nCoverage = FPs above threshold / 5\nBalance = 1 - spread / worst spread\nAlignment = do teachers & leaders agree?\nC1 = 10 x Teacher(Cov x Bal x Align) + 10 x Leader(...)"),
        ("C2 Practice", "25", C2_CLR,
         "Build 5x4 depth grid, remove impossible cells\nAverage depth per FP, normalise to 0-1\nDepth Effectiveness = coverage x avg normalised depth\nC2 = 25 x Depth Effectiveness"),
        ("C3 Capacity", "20", C3_CLR,
         "4 dimensions: higher-order, harder problems, revision, consistency\nEach 0-1. CPI = average of 4 dimensions\nBranch CPI = average across all teachers\nC3 = 20 x Branch CPI"),
        ("C4 System", "20", C4_CLR,
         "Agree/SA = 1, Disagree/SD = 0\nEnablement = avg(A1, A2, A3, A4)\nRoutine = avg(B1, B2, B3, B4, B5)\nC4 = 10 x Enablement + 10 x Routine"),
    ]
    for i, (name, mx, clr, desc) in enumerate(formulas):
        col = i % 2
        row = i // 2
        left = Inches(0.5 + col * 6.4)
        top = Inches(1.4 + row * 2.9)
        bx = add_shape(slide, left, top, Inches(6.0), Inches(2.6))
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Pt(5), Inches(2.6))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = clr
        stripe.line.fill.background()
        tf = set_text(bx, f"    {name} (max {mx})", size=15, bold=True, color=clr)
        for line in desc.split("\n"):
            add_para(tf, f"    {line}", size=11, color=TEXT_CLR, space_before=4)

    # ── Slide 16: Key Takeaways ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title_bar(slide, "Key Takeaways")

    takeaways = [
        "SRI gives each school a single 0-100 score built from 5 independent, clearly-defined areas.",
        "No single area dominates — a school must be strong across Intent, Practice, Capacity, and System.",
        "C1 (Intent) multiplies three factors together: miss one and the whole score drops. Both roles must be present and aligned.",
        "C2 (Practice) rewards depth over surface awareness — 'doing beats knowing'. Impossible-to-observe combinations are filtered out.",
        "C3 (Capacity) uses an adaptive test with 4 thinking dimensions. It's personalised — harder questions for stronger thinkers.",
        "C4 (System) simplifies a 4-point survey into Yes/No, then averages across teacher and leader areas.",
        "Every formula step has a clear, plain-English interpretation. No black boxes.",
        "With ~65/85 on average, DES branches show strong readiness with room for deeper practice and broader participation.",
    ]
    for i, t in enumerate(takeaways):
        top = Inches(1.3 + i * 0.75)
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), top + Inches(0.06), Inches(0.42), Inches(0.42))
        circ.fill.solid()
        circ.fill.fore_color.rgb = ACCENT
        circ.line.fill.background()
        set_text(circ, str(i + 1), size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
        bx = add_shape(slide, Inches(1.4), top, Inches(11.4), Inches(0.56))
        set_text(bx, f"  {t}", size=13, color=TEXT_CLR)

    # ── Slide 17: Thank You ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PRIMARY)
    box = add_shape(slide, Inches(1), Inches(1.5), Inches(11.3), Inches(4.5), PRIMARY)
    box.line.fill.background()
    tf = set_text(box, "Thank You", size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_para(tf, "", size=14, space_before=16)
    add_para(tf, "Every number in SRI has a clear, traceable meaning.", size=20, color=RGBColor(0x93, 0xC5, 0xFD), space_before=12)
    add_para(tf, "No black boxes. No hidden assumptions.", size=20, color=RGBColor(0x93, 0xC5, 0xFD), space_before=4)
    add_para(tf, "", size=14, space_before=24)
    add_para(tf, "Myelin NEP Readiness Pilot", size=16, color=RGBColor(0x93, 0xC5, 0xFD), space_before=8)
    add_para(tf, "Deccan Education Society  |  February 2026", size=14, color=RGBColor(0x93, 0xC5, 0xFD), space_before=8)

    path = DOCS / "SRI_Methodology_Explained.pptx"
    prs.save(str(path))
    print(f"PPTX -> {path}  ({len(prs.slides)} slides)")
    return path


# ═══════════════════════════════════════════════════════════════
if __name__ == "__main__":
    html_path = generate_html()
    pptx_path = generate_pptx()
    print(f"\nDone. Open in browser for PDF: {html_path}")
    print(f"Presentation ready: {pptx_path}")
