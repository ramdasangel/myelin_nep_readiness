#!/usr/bin/env python3
"""Generate PPTX for Functional Steps: User Task Mapping & Daily Logging."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

OUTPUT = Path(__file__).resolve().parent.parent / "docs" / "functional_steps_task_mapping_logging.pptx"

# Colors
PRIMARY = RGBColor(0x1A, 0x36, 0x5D)
ACCENT = RGBColor(0x2B, 0x6C, 0xB0)
GREEN = RGBColor(0x38, 0xA1, 0x69)
AMBER = RGBColor(0xD6, 0x9E, 0x2E)
RED = RGBColor(0xE5, 0x3E, 0x3E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF7, 0xFA, 0xFC)
MUTED = RGBColor(0x71, 0x80, 0x96)
TEXT_CLR = RGBColor(0x2D, 0x37, 0x48)
FP_COLORS = {
    'FP-1': RGBColor(0x63, 0x66, 0xF1),
    'FP-2': RGBColor(0xEC, 0x48, 0x99),
    'FP-3': RGBColor(0xF5, 0x9E, 0x0B),
    'FP-4': RGBColor(0x10, 0xB9, 0x81),
    'FP-5': RGBColor(0x3B, 0x82, 0xF6),
}

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color=LIGHT_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color=WHITE, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def set_text(shape, text, size=12, bold=False, color=TEXT_CLR, alignment=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return tf

def add_para(tf, text, size=12, bold=False, color=TEXT_CLR, space_before=4):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.space_before = Pt(space_before)
    return p

# ─── SLIDE 1: Title ───
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, PRIMARY)
shape = add_shape(slide, Inches(1), Inches(1.5), Inches(11.3), Inches(4.5), PRIMARY)
shape.line.fill.background()
tf = set_text(shape, "Functional Steps", size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_para(tf, "User Task Mapping & Daily Logging", size=28, color=WHITE, space_before=8)
add_para(tf, "", size=16, color=WHITE, space_before=20)
add_para(tf, "Stage 5 — Micro-Intervention Practice", size=18, color=RGBColor(0x93, 0xC5, 0xFD), space_before=16)
add_para(tf, "Myelin NEP Readiness Pilot  |  Deccan Education Society", size=14, color=RGBColor(0x93, 0xC5, 0xFD), space_before=8)
add_para(tf, "2026-02-18", size=14, color=RGBColor(0x93, 0xC5, 0xFD), space_before=8)

# ─── SLIDE 2: Overview ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
s = add_shape(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8), PRIMARY)
set_text(s, "  Overview", size=24, bold=True, color=WHITE)
# Description
s = add_shape(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(1.2))
tf = set_text(s, "Stage 5 enables teachers and school leaders to select daily micro-intervention tasks aligned to NEP 2020's Foundational Principles and track classroom practice through a 21-day daily logging mechanism.", size=14, color=TEXT_CLR)
add_para(tf, "User Roles: Teacher (~70%) and Leader (~30%)  |  Languages: English + Marathi  |  Duration: 21-day practice window", size=12, color=MUTED, space_before=12)

# KPI boxes
kpis = [("12", "Tasks Available", ACCENT), ("5", "Foundational\nPrinciples", GREEN), ("21", "Day Practice\nWindow", AMBER), ("2", "Languages\n(EN + MR)", ACCENT), ("2", "Roles\n(Teacher + Leader)", ACCENT)]
for i, (val, label, clr) in enumerate(kpis):
    left = Inches(0.5 + i * 2.5)
    box = add_shape(slide, left, Inches(2.8), Inches(2.2), Inches(1.6))
    # top border line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(2.8), Inches(2.2), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = clr
    line.line.fill.background()
    tf = set_text(box, val, size=36, bold=True, color=clr, alignment=PP_ALIGN.CENTER)
    add_para(tf, label, size=11, color=MUTED, space_before=6)
    tf.paragraphs[-1].alignment = PP_ALIGN.CENTER

# 3-phase overview boxes
phases = [
    ("Phase 1: Task Mapping", "User selects 3-12 tasks\nfrom the task library", FP_COLORS['FP-1']),
    ("Phase 2: Daily Logging", "Daily check-in for 21 days\nMark practiced tasks + reflect", GREEN),
    ("Phase 3: Progress Tracking", "Dashboard shows engagement,\ncompletion, and consistency", FP_COLORS['FP-5']),
]
for i, (title, desc, clr) in enumerate(phases):
    left = Inches(0.5 + i * 4.2)
    box = add_shape(slide, left, Inches(4.8), Inches(3.8), Inches(2.3))
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(4.8), Inches(3.8), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = clr
    line.line.fill.background()
    tf = set_text(box, title, size=16, bold=True, color=PRIMARY, alignment=PP_ALIGN.CENTER)
    add_para(tf, desc, size=12, color=MUTED, space_before=10)
    tf.paragraphs[-1].alignment = PP_ALIGN.CENTER

# ─── SLIDE 3: Foundational Principles ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
s = add_shape(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8), PRIMARY)
set_text(s, "  NEP 2020 Foundational Principles", size=24, bold=True, color=WHITE)

fp_data = [
    ("FP-1", "Every Child\nis Unique", "Recognizing individual\ndifferences, inclusive\nobservation", "T01, T02, T11"),
    ("FP-2", "Holistic &\nExperiential Learning", "Learning by doing,\nconnecting to\nreal life", "T03, T04"),
    ("FP-3", "Teacher as Reflective\nPractitioner", "Self-reflection,\ncontinuous\nimprovement", "T05, T06"),
    ("FP-4", "Assessment\nfor Learning", "Formative assessment,\nunderstanding\nover marks", "T07, T08, T12"),
    ("FP-5", "Collaboration\n& Community", "Teacher-parent-\ncommunity\npartnership", "T09, T10"),
]
for i, (code, name, desc, tasks) in enumerate(fp_data):
    left = Inches(0.5 + i * 2.5)
    box = add_shape(slide, left, Inches(1.4), Inches(2.2), Inches(5.5))
    clr = FP_COLORS[code]
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(1.4), Inches(2.2), Pt(5))
    line.fill.solid()
    line.fill.fore_color.rgb = clr
    line.line.fill.background()
    tf = set_text(box, code, size=22, bold=True, color=clr, alignment=PP_ALIGN.CENTER)
    add_para(tf, name, size=14, bold=True, color=PRIMARY, space_before=12)
    tf.paragraphs[-1].alignment = PP_ALIGN.CENTER
    add_para(tf, desc, size=11, color=MUTED, space_before=14)
    tf.paragraphs[-1].alignment = PP_ALIGN.CENTER
    add_para(tf, "Tasks:", size=10, bold=True, color=PRIMARY, space_before=16)
    tf.paragraphs[-1].alignment = PP_ALIGN.CENTER
    add_para(tf, tasks, size=11, color=clr, space_before=4)
    tf.paragraphs[-1].alignment = PP_ALIGN.CENTER

# ─── SLIDE 4: 12 Tasks (Part 1 - T01 to T06) ───
tasks_data = [
    ("T01", "Notice One Learner", "FP-1", "D1->D2", "Quietly watch one student - notice how they sit, listen, respond, or try"),
    ("T02", "One Adjustment", "FP-1", "D2->D3", "Change one small thing for a student - speak slower, simpler example, or pair them"),
    ("T03", "Change One Example", "FP-2", "D1->D2", "Replace textbook example with something from daily life students know"),
    ("T04", "Ask a Why/How", "FP-2", "D2->D3", "Ask at least one 'Why?' or 'How did you get this answer?'"),
    ("T05", "End-of-Class Reflection", "FP-3", "D1->D2", "After class, think: What went well? What didn't work?"),
    ("T06", "Try One Small Change", "FP-3", "D2->D3", "In next class, change one thing based on what you noticed"),
    ("T07", "Quick Check", "FP-4", "D1->D2", "Ask 2-3 students to explain the answer in their own words"),
    ("T08", "Spot One Pattern", "FP-4", "D2->D3", "Notice if many students make the same mistake or confusion"),
    ("T09", "Teacher Touchpoint", "FP-5", "D1->D2", "Tell another teacher about one thing that happened in your class"),
    ("T10", "Parent Signal", "FP-5", "D1->D2", "Share one good learning observation about a child with a parent"),
    ("T11", "Student Voice", "FP-1", "Cross-FP", "Ask one student: 'What helped you learn today?'"),
    ("T12", "Pause & Name", "FP-4", "Cross-FP", "Pause once and name what students are doing well in learning"),
]

for part in range(2):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    start = part * 6
    end = start + 6
    title_text = f"  The 12 Micro-Intervention Tasks ({start+1}-{end})"
    s = add_shape(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8), PRIMARY)
    set_text(s, title_text, size=24, bold=True, color=WHITE)

    for i, (code, name, fp, depth, desc) in enumerate(tasks_data[start:end]):
        row = i
        top = Inches(1.4 + row * 0.95)
        # Task card
        box = add_shape(slide, Inches(0.5), top, Inches(12.3), Inches(0.85))
        clr = FP_COLORS[fp]
        # Left color stripe
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), top, Pt(5), Inches(0.85))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = clr
        stripe.line.fill.background()
        # Code
        tf = set_text(box, f"    {code}", size=16, bold=True, color=clr)
        # Name + FP + Depth
        run_text = f"  {name}     [{fp}]  [{depth}]"
        add_para(tf, f"    {name}        {fp}  |  {depth}", size=12, bold=False, color=PRIMARY, space_before=2)
        add_para(tf, f"    {desc}", size=11, color=MUTED, space_before=2)

# ─── SLIDE 6: Depth Levels ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
s = add_shape(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8), PRIMARY)
set_text(s, "  Depth Levels Explained", size=24, bold=True, color=WHITE)

depths = [
    ("D1 -> D2", "Procedural -> Responsive", "Entry-level: start by noticing patterns, then respond to what you observe.", "60% of selections", ACCENT),
    ("D2 -> D3", "Responsive -> Diagnostic", "Growth: adjust responses based on deeper evidence and diagnosis.", "30% of selections", GREEN),
    ("Cross-FP", "Multi-Principle Integration", "Advanced: tasks that bridge multiple foundational principles.", "10% of selections", AMBER),
]
for i, (code, label, desc, pct, clr) in enumerate(depths):
    top = Inches(1.5 + i * 1.8)
    box = add_shape(slide, Inches(0.5), top, Inches(12.3), Inches(1.5))
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), top, Pt(6), Inches(1.5))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = clr
    stripe.line.fill.background()
    tf = set_text(box, f"     {code}", size=20, bold=True, color=clr)
    add_para(tf, f"     {label}", size=14, bold=True, color=PRIMARY, space_before=6)
    add_para(tf, f"     {desc}", size=12, color=MUTED, space_before=6)
    add_para(tf, f"     {pct}", size=11, bold=True, color=clr, space_before=6)

# ─── SLIDE 7: Phase 1 - Task Mapping ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
s = add_shape(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8), FP_COLORS['FP-1'])
set_text(s, "  Phase 1: Task Mapping (Enrollment)", size=24, bold=True, color=WHITE)

steps_p1 = [
    "User logs in and navigates to the Micro-Intervention module",
    "System displays all 12 tasks organized by FP — each with name, description, FP tag, depth tag. Bilingual toggle: EN | Marathi",
    "User selects 3-12 tasks. Selection highlights FP coverage and depth distribution visually",
    'User confirms selection -> "Start My Practice"',
    "System saves to UserTaskMapping: userTempId, SelectedTasks[], createdAt",
    "User sees confirmation screen with selected tasks and practice start date",
]
for i, step in enumerate(steps_p1):
    top = Inches(1.4 + i * 0.85)
    # Number circle
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), top + Inches(0.1), Inches(0.5), Inches(0.5))
    circ.fill.solid()
    circ.fill.fore_color.rgb = FP_COLORS['FP-1']
    circ.line.fill.background()
    set_text(circ, str(i+1), size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    # Step text
    box = add_shape(slide, Inches(1.4), top, Inches(11.4), Inches(0.65))
    set_text(box, step, size=13, color=TEXT_CLR)

# Business rules
box = add_shape(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.7), RGBColor(0xF0, 0xFF, 0xF4))
set_text(box, "  Business Rules: Min 3 tasks, max 12. User can re-select anytime. Same task library for Teachers and Leaders.", size=11, bold=False, color=GREEN)

# ─── SLIDE 8: Phase 2 - Daily Logging ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
s = add_shape(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8), GREEN)
set_text(s, "  Phase 2: Daily Logging (21-Day Practice Window)", size=24, bold=True, color=WHITE)

steps_p2 = [
    'User opens "Daily Practice Log" screen. Current date shown prominently',
    "System displays ONLY the user's selected tasks (not all 12). Each task as a card: checkbox + optional comment field",
    "User checks off completed tasks: true (practiced) / false (not practiced)",
    "User optionally adds a reflection comment (free text)",
    "System auto-saves on each interaction (no submit button)",
    "System creates UserDailyProgress record: UserTempId, SubmitDate, TasksProgress[{taskId, isChecked, comment}]",
    "User sees visual confirmation - checkmarks turn green",
]
for i, step in enumerate(steps_p2):
    top = Inches(1.4 + i * 0.75)
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), top + Inches(0.05), Inches(0.5), Inches(0.5))
    circ.fill.solid()
    circ.fill.fore_color.rgb = GREEN
    circ.line.fill.background()
    set_text(circ, str(i+1), size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    box = add_shape(slide, Inches(1.4), top, Inches(11.4), Inches(0.58))
    set_text(box, step, size=13, color=TEXT_CLR)

box = add_shape(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.7), RGBColor(0xFF, 0xFB, 0xEB))
set_text(box, "  Business Rules: One log per user per day. All selected tasks appear. Comments optional. Append-only (no edits/deletes).", size=11, color=AMBER)

# ─── SLIDE 9: Phase 3 - Progress Tracking ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
s = add_shape(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8), FP_COLORS['FP-5'])
set_text(s, "  Phase 3: Progress Tracking (Dashboard)", size=24, bold=True, color=WHITE)

steps_p3 = [
    'User navigates to "My Progress" screen',
    "System displays 21-day progress grid: Rows = tasks, Columns = dates. Cell states: Done | Done+comment | Not done | Future",
    "System calculates per-user metrics: completion rate, consistency %, comment rate",
    "Admin dashboard aggregates across all users: KPIs, charts (daily timeline, task completion, FP radar), filterable user table",
]
for i, step in enumerate(steps_p3):
    top = Inches(1.4 + i * 1.0)
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), top + Inches(0.1), Inches(0.5), Inches(0.5))
    circ.fill.solid()
    circ.fill.fore_color.rgb = FP_COLORS['FP-5']
    circ.line.fill.background()
    set_text(circ, str(i+1), size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    box = add_shape(slide, Inches(1.4), top, Inches(11.4), Inches(0.75))
    set_text(box, step, size=14, color=TEXT_CLR)

# Metrics table
s = add_shape(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.5))
set_text(s, "  Key Metrics", size=16, bold=True, color=PRIMARY)

metrics = [
    ("Completion Rate", "totalChecked / (totalChecked + totalUnchecked) x 100", "Per user + global"),
    ("Consistency", "daysLogged / min(daysSinceMapping, 21) x 100", "Per user"),
    ("Comment Rate", "totalComments / totalLogEntries x 100", "Per user + global"),
]
for i, (name, formula, scope) in enumerate(metrics):
    top = Inches(5.8 + i * 0.5)
    box = add_shape(slide, Inches(0.5), top, Inches(3.5), Inches(0.42))
    set_text(box, f"  {name}", size=12, bold=True, color=PRIMARY)
    box = add_shape(slide, Inches(4.1), top, Inches(6.0), Inches(0.42))
    set_text(box, formula, size=11, color=MUTED)
    box = add_shape(slide, Inches(10.2), top, Inches(2.6), Inches(0.42))
    set_text(box, scope, size=11, color=ACCENT, alignment=PP_ALIGN.CENTER)

# ─── SLIDE 10: Data Flow Diagram ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
s = add_shape(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8), PRIMARY)
set_text(s, "  End-to-End Data Flow", size=24, bold=True, color=WHITE)

# User journey row
flow_items = [
    ("LOGIN", "User\nauthenticates"),
    ("SELECT TASKS", "Choose\n3-12 tasks"),
    ("DAILY LOG", "Check off\n+ reflect"),
    ("DASHBOARD", "View\nprogress"),
]
for i, (title, desc) in enumerate(flow_items):
    left = Inches(0.5 + i * 3.3)
    clr = ACCENT if i < 3 else AMBER
    box = add_shape(slide, left, Inches(1.5), Inches(2.5), Inches(1.3), RGBColor(0xEB, 0xF4, 0xFF), clr)
    tf = set_text(box, title, size=14, bold=True, color=PRIMARY, alignment=PP_ALIGN.CENTER)
    add_para(tf, desc, size=11, color=MUTED, space_before=6)
    tf.paragraphs[-1].alignment = PP_ALIGN.CENTER
    if i < 3:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left + Inches(2.6), Inches(1.9), Inches(0.6), Inches(0.4))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ACCENT
        arrow.line.fill.background()

# Down arrow
arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.2), Inches(3.0), Inches(0.5), Inches(0.6))
arrow.fill.solid()
arrow.fill.fore_color.rgb = ACCENT
arrow.line.fill.background()

# Database row
db_items = [
    ("UserTaskMapping", "userTempId\nSelectedTasks[]\ncreatedAt"),
    ("UserDailyProgress", "UserTempId, SubmitDate\nTasksProgress[]\n{taskId, isChecked, comment}"),
    ("Computed Metrics", "Completion Rate\nConsistency %\nComment Rate"),
]
for i, (title, desc) in enumerate(db_items):
    left = Inches(0.5 + i * 4.3)
    bg_clr = RGBColor(0xF0, 0xFF, 0xF4) if i < 2 else RGBColor(0xFF, 0xFB, 0xEB)
    bd_clr = GREEN if i < 2 else AMBER
    box = add_shape(slide, left, Inches(3.9), Inches(3.8), Inches(1.8), bg_clr, bd_clr)
    tf = set_text(box, title, size=14, bold=True, color=PRIMARY, alignment=PP_ALIGN.CENTER)
    add_para(tf, desc, size=11, color=MUTED, space_before=8)
    tf.paragraphs[-1].alignment = PP_ALIGN.CENTER
    if i < 2:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left + Inches(3.9), Inches(4.5), Inches(0.35), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GREEN
        arrow.line.fill.background()

# ─── SLIDE 11: User Scenarios ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
s = add_shape(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8), PRIMARY)
set_text(s, "  Typical User Scenarios", size=24, bold=True, color=WHITE)

scenarios = [
    ("Scenario A: Engaged Teacher", "Teacher", GREEN,
     "Feb 3: Maps 5 tasks (one per FP, D1->D2)\nFeb 4-17: Logs 11 of 14 days, 4 comments",
     "Consistency: 79%  |  Completion: 85%  |  Comments: 7%"),
    ("Scenario B: Selective Leader", "Leader", AMBER,
     "Feb 4: Maps 3 growth tasks (D2->D3)\nFeb 5-17: Logs 6 of 13 days",
     "Consistency: 46%  |  Completion: 72%  |  Comments: 0%"),
    ("Scenario C: Drop-Off Teacher", "Teacher", RED,
     "Feb 3: Maps 7 tasks across all FPs\nFeb 5-7: Logs 3 days, then stops",
     "Consistency: 20%  |  Completion: 95%  |  Comments: 15%"),
]
for i, (title, role, clr, timeline, result) in enumerate(scenarios):
    left = Inches(0.5 + i * 4.2)
    box = add_shape(slide, left, Inches(1.4), Inches(3.8), Inches(5.5))
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(1.4), Inches(3.8), Pt(5))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = clr
    stripe.line.fill.background()
    tf = set_text(box, title, size=15, bold=True, color=PRIMARY, alignment=PP_ALIGN.CENTER)
    add_para(tf, role, size=12, bold=True, color=clr, space_before=8)
    tf.paragraphs[-1].alignment = PP_ALIGN.CENTER
    add_para(tf, "", size=8, space_before=10)
    add_para(tf, timeline, size=11, color=TEXT_CLR, space_before=8)
    add_para(tf, "", size=8, space_before=14)
    add_para(tf, result, size=11, bold=True, color=clr, space_before=10)
    tf.paragraphs[-1].alignment = PP_ALIGN.CENTER

# ─── SLIDE 12: Design Principles ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
s = add_shape(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8), PRIMARY)
set_text(s, "  Design Principles", size=24, bold=True, color=WHITE)

principles = [
    ("Trust-Based Reflection", "Self-reported check-ins, not surveillance.\nNo external validation required.", FP_COLORS['FP-1']),
    ("Minimal Friction", "Auto-save on interaction. No submit button.\nCard-based mobile UI.", GREEN),
    ("Bilingual From Day 1", "English + Marathi toggle on all screens.\nTasks, descriptions, UI labels.", FP_COLORS['FP-3']),
    ("Mobile-First", "Cards on mobile, tables on desktop.\nResponsive layout throughout.", FP_COLORS['FP-5']),
    ("Append-Only Data", "No edits or deletes.\nAll log entries are permanent and immutable.", RED),
    ("Role-Agnostic Tasks", "Same 12-task library for Teachers and Leaders.\nContext differs, tasks don't.", AMBER),
]
for i, (title, desc, clr) in enumerate(principles):
    col = i % 2
    row = i // 2
    left = Inches(0.5 + col * 6.3)
    top = Inches(1.4 + row * 1.9)
    box = add_shape(slide, left, top, Inches(5.8), Inches(1.6))
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Pt(5), Inches(1.6))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = clr
    stripe.line.fill.background()
    tf = set_text(box, f"    {title}", size=15, bold=True, color=PRIMARY)
    add_para(tf, f"    {desc}", size=12, color=MUTED, space_before=8)

# ─── SLIDE 13: Current Stats ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
s = add_shape(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8), PRIMARY)
set_text(s, "  Current Production Statistics (2026-02-18)", size=24, bold=True, color=WHITE)

stats = [
    ("124", "Total Users\n(DES)", ACCENT),
    ("119", "Active\n(>=1 Log)", GREEN),
    ("5", "Zero-Log\nUsers", RED),
    ("5.7", "Avg Days\nLogged", AMBER),
    ("~88%", "Completion\nRate", GREEN),
    ("2,994", "Total Log\nEntries", ACCENT),
]
for i, (val, label, clr) in enumerate(stats):
    col = i % 3
    row = i // 3
    left = Inches(0.5 + col * 4.2)
    top = Inches(1.5 + row * 2.8)
    box = add_shape(slide, left, Inches(1.5 + row * 2.8), Inches(3.8), Inches(2.3))
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(3.8), Pt(5))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = clr
    stripe.line.fill.background()
    tf = set_text(box, val, size=48, bold=True, color=clr, alignment=PP_ALIGN.CENTER)
    add_para(tf, label, size=13, color=MUTED, space_before=8)
    tf.paragraphs[-1].alignment = PP_ALIGN.CENTER

# Save
prs.save(str(OUTPUT))
print(f"PPTX saved to {OUTPUT}")
print(f"  {len(prs.slides)} slides")
